import json
import os
import re
import base64
import hashlib
import shutil
import subprocess
from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path, PurePosixPath
from typing import Any, Optional
from uuid import uuid4

import httpx
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from sqlalchemy import delete, func, or_, select, text
from sqlalchemy.orm import Session

from app import models
from app.config import settings
from app.skill_manifest import SKILL_BY_ID, SkillDefinition, list_builtin_skills, match_skill_by_keywords

TENCENT_MEETING_SCRIPT = Path.home() / ".codex" / "skills" / "tencent-meeting-mcp" / "scripts" / "tencent_meeting.py"


class TencentMinutesUnavailableError(RuntimeError):
    pass


SUPPORTED_PROJECT_EXTS = {".pdf", ".docx", ".xlsx", ".txt", ".md", ".pptx", ".png", ".jpg", ".jpeg"}
SUPPORTED_INBOX_EXTS = SUPPORTED_PROJECT_EXTS | {".zip", ".rar", ".7z", ".dwg", ".skp", ".rvt", ".csv", ".webp", ".mp3", ".m4a", ".wav", ".mp4"}
SUPPORTED_KNOWLEDGE_EXTS = {".md", ".txt", ".pdf", ".docx", ".xlsx", ".pptx", ".png", ".jpg", ".jpeg"}
MAX_BROWSER_UPLOAD_SIZE = 25 * 1024 * 1024
DEFAULT_VAULT_EXCLUDE_DIRS = {
    ".git",
    ".obsidian",
    ".venv",
    ".venv-funasr",
    "__pycache__",
    "node_modules",
    "dist",
    ".cache",
}
DEFAULT_VAULT_EXCLUDE_EXTS = {
    ".mp4",
    ".mov",
    ".avi",
    ".zip",
    ".rar",
    ".7z",
    ".psd",
    ".dwg",
    ".bak",
}
DEFAULT_VAULT_MAX_FILE_SIZE = 30 * 1024 * 1024
KNOWLEDGE_OVERVIEW_TRIGGERS = ("知识库里有什么", "知识库有什么", "有什么资料", "资料概览", "知识库概览", "总结知识库")

DIGITAL_EMPLOYEES = [
    ("AI项目经理", "项目负责人 / PM", ["任务拆解", "进度协调", "风险跟踪"], "PM"),
    ("AI总图助理", "总图负责人", ["强排", "流线", "指标复核"], "MP"),
    ("AI户型助理", "户型负责人", ["产品定位", "户型推演", "得房率"], "UN"),
    ("AI立面助理", "立面负责人", ["风格研究", "材料策略", "比例控制"], "FC"),
    ("AI景观助理", "景观接口", ["景观节点", "归家动线", "界面协同"], "LS"),
    ("AI室内助理", "室内接口", ["大堂", "会所", "精装氛围"], "IN"),
    ("AI汇报助理", "汇报助理", ["PPT目录", "汇报文案", "图像提示词"], "RP"),
    ("AI资料管理员", "资料管理员", ["文件整理", "知识库引用", "版本记录"], "KB"),
]


def parse_tencent_meeting_result(text: str) -> dict[str, str]:
    text = text.replace("\\n", "\n")
    code_match = re.search(r"(?:会议号|meeting_code|meeting code)\s*[：:]\s*([0-9\s-]{9,})", text, re.I)
    url_match = re.search(r"https://meeting\.tencent\.com/[^\\\s，。)）]+", text)
    trace_match = re.search(r"(?:X-Tc-Trace|rpcUuid)\s*[：:]\s*([^\s，。]+)", text, re.I)
    meeting_id_match = re.search(r"(?:meeting_id|会议ID|会议id)\s*[：:]\s*([0-9]+)", text, re.I)
    return {
        "meeting_code": re.sub(r"\D", "", code_match.group(1)) if code_match else "",
        "join_url": url_match.group(0) if url_match else "",
        "trace": trace_match.group(1) if trace_match else "",
        "meeting_id": meeting_id_match.group(1) if meeting_id_match else "",
        "raw": text,
    }


def _load_tencent_json_payload(text: str) -> Any:
    try:
        outer = json.loads(text)
    except json.JSONDecodeError:
        return None
    body = outer.get("body") if isinstance(outer, dict) else None
    if isinstance(body, str):
        try:
            return json.loads(body)
        except json.JSONDecodeError:
            return outer
    return outer


def _tencent_response_error(text: str) -> str:
    try:
        payload = json.loads(text)
    except json.JSONDecodeError:
        payload = None
    if not isinstance(payload, dict):
        return ""
    status_code = int(payload.get("status_code") or 200)
    if status_code < 400:
        return ""
    body = payload.get("body")
    body_payload: Any = {}
    if isinstance(body, str):
        try:
            body_payload = json.loads(body)
        except json.JSONDecodeError:
            body_payload = {}
    message = ""
    if isinstance(body_payload, dict):
        error_info = body_payload.get("error_info")
        if isinstance(error_info, dict):
            message = str(error_info.get("message") or "")
    trace = ""
    headers = payload.get("headers")
    if isinstance(headers, dict):
        trace = str(headers.get("X-Tc-Trace") or headers.get("rpcUuid") or "")
    return "；".join(part for part in [message or f"腾讯会议接口返回 {status_code}", f"trace={trace}" if trace else ""] if part)


def extract_tencent_transcript_text(text: str) -> str:
    payload = _load_tencent_json_payload(text)
    body = payload.get("body") if isinstance(payload, dict) else None
    if isinstance(body, str):
        try:
            body_payload = json.loads(body)
        except json.JSONDecodeError:
            return text
    else:
        body_payload = payload
    if not isinstance(body_payload, dict):
        return text
    minutes = body_payload.get("minutes")
    if not isinstance(minutes, dict):
        return text
    lines: list[str] = []
    for paragraph in minutes.get("paragraphs", []) or []:
        if not isinstance(paragraph, dict):
            continue
        speaker = paragraph.get("speaker")
        speaker_name = ""
        if isinstance(speaker, dict):
            speaker_name = str(speaker.get("user_name") or "").strip()
        words: list[str] = []
        for sentence in paragraph.get("sentences", []) or []:
            if not isinstance(sentence, dict):
                continue
            for word in sentence.get("words", []) or []:
                if isinstance(word, dict):
                    words.append(str(word.get("text") or ""))
        content = "".join(words).strip()
        if content:
            lines.append(f"{speaker_name}：{content}" if speaker_name else content)
    return "\n".join(lines)


def extract_tencent_smart_minutes_text(text: str) -> str:
    payload = _load_tencent_json_payload(text)
    body = payload.get("body") if isinstance(payload, dict) else None
    if isinstance(body, str):
        try:
            body_payload = json.loads(body)
        except json.JSONDecodeError:
            return text
    else:
        body_payload = payload
    if not isinstance(body_payload, dict):
        return text
    meeting_minute = body_payload.get("meeting_minute")
    if isinstance(meeting_minute, dict):
        minute = str(meeting_minute.get("minute") or "").strip()
        if minute:
            return minute
    return text


def parse_tencent_record_result(text: str) -> dict[str, Any]:
    text = text.replace("\\n", "\n")
    payload = _load_tencent_json_payload(text)
    records: list[dict[str, str]] = []
    if isinstance(payload, dict):
        for meeting_record in payload.get("record_meetings", []) or []:
            if not isinstance(meeting_record, dict):
                continue
            meeting_record_id = str(meeting_record.get("meeting_record_id") or "")
            for record_file in meeting_record.get("record_files", []) or []:
                if not isinstance(record_file, dict):
                    continue
                record_file_id = str(record_file.get("record_file_id") or "")
                if record_file_id:
                    records.append({"record_file_id": record_file_id, "meeting_record_id": meeting_record_id})

    record_file_match = re.search(r"(?:record_file_id|recordFileId|录制文件ID|录制文件id)\s*[：:]\s*([0-9]+)", text, re.I)
    meeting_record_match = re.search(r"(?:meeting_record_id|meetingRecordId|会议录制ID|会议录制id)\s*[：:]\s*([0-9]+)", text, re.I)
    download_match = re.search(r"https?://[^\\\s，。)）]+", text)
    first_record = records[0] if records else {}
    return {
        "record_file_id": first_record.get("record_file_id") or (record_file_match.group(1) if record_file_match else ""),
        "meeting_record_id": first_record.get("meeting_record_id") or (meeting_record_match.group(1) if meeting_record_match else ""),
        "download_url": download_match.group(0) if download_match else "",
        "records": records,
        "raw": text,
    }


def parse_tencent_record_address_result(text: str) -> dict[str, str]:
    payload = _load_tencent_json_payload(text)
    view_address = ""
    if isinstance(payload, dict):
        for record_file in payload.get("record_files", []) or []:
            if isinstance(record_file, dict) and record_file.get("view_address"):
                view_address = str(record_file["view_address"])
                break
    outer = json.loads(text) if text.strip().startswith("{") else {}
    headers = outer.get("headers") if isinstance(outer, dict) else {}
    return {
        "view_address": view_address,
        "trace": str(headers.get("X-Tc-Trace") or "") if isinstance(headers, dict) else "",
        "rpc_uuid": str(headers.get("rpcUuid") or "") if isinstance(headers, dict) else "",
    }


def extract_tencent_meeting_metadata(meeting: models.Meeting) -> dict[str, str]:
    legacy = parse_tencent_meeting_result("\n".join([meeting.recording_url or "", meeting.agenda or ""]))
    return {
        **legacy,
        "meeting_code": meeting.tencent_meeting_code or legacy["meeting_code"],
        "meeting_id": meeting.tencent_meeting_id or legacy["meeting_id"],
        "join_url": meeting.tencent_join_url or legacy["join_url"],
    }


def extract_tencent_transcript_paragraph_ids(text: str) -> list[str]:
    payload = _load_tencent_json_payload(text)
    found: list[str] = []

    def visit(value: Any) -> None:
        if isinstance(value, dict):
            pid = value.get("pid") or value.get("paragraph_id")
            if pid is not None and str(pid) not in found:
                found.append(str(pid))
            for child in value.values():
                visit(child)
        elif isinstance(value, list):
            for child in value:
                visit(child)

    visit(payload)
    return found


def fetch_complete_tencent_transcript(record_file_id: str) -> str:
    try:
        paragraph_text = call_tencent_meeting_tool("get_transcripts_paragraphs", {"record_file_id": record_file_id})
        paragraph_ids = extract_tencent_transcript_paragraph_ids(paragraph_text)
    except Exception:
        paragraph_ids = []
    if not paragraph_ids:
        paragraph_ids = ["0"]
    lines: list[str] = []
    seen: set[str] = set()
    for pid in paragraph_ids:
        detail = call_tencent_meeting_tool(
            "get_transcripts_details",
            {"record_file_id": record_file_id, "pid": pid, "limit": "50"},
        )
        for line in extract_tencent_transcript_text(detail).splitlines():
            normalized = line.strip()
            if normalized and normalized not in seen:
                seen.add(normalized)
                lines.append(normalized)
    return "\n".join(lines)


def call_tencent_meeting_tool(name: str, arguments: dict[str, Any], timeout: int = 60) -> str:
    token = (settings.tencent_meeting_token or os.environ.get("TENCENT_MEETING_TOKEN", "")).strip()
    if not token:
        raise RuntimeError("TENCENT_MEETING_TOKEN 未配置")
    if not TENCENT_MEETING_SCRIPT.exists():
        raise RuntimeError("腾讯会议 skill 未安装")
    payload = {
        "name": name,
        "arguments": {
            **arguments,
            "_client_info": {"os": "macos", "agent": "codex-rom-ai", "model": "gpt-5-codex"},
        },
    }
    result = subprocess.run(
        ["python3", str(TENCENT_MEETING_SCRIPT), "tools/call", json.dumps(payload, ensure_ascii=False)],
        text=True,
        capture_output=True,
        timeout=timeout,
        env={**os.environ, "TENCENT_MEETING_TOKEN": token},
    )
    output = (result.stdout or "") + ("\n" + result.stderr if result.stderr else "")
    response_error = _tencent_response_error(output)
    if response_error:
        raise RuntimeError(response_error)
    if result.returncode != 0 or "[错误]" in output:
        raise RuntimeError(output.strip() or "腾讯会议创建失败")
    return output


def schedule_tencent_meeting(subject: str, start_time: str, end_time: str) -> dict[str, str]:
    output = call_tencent_meeting_tool(
        "schedule_meeting",
        {
            "subject": subject,
            "start_time": start_time,
            "end_time": end_time,
            "time_zone": "Asia/Shanghai",
        },
    )
    return parse_tencent_meeting_result(output)


def build_tencent_records_query(meeting: models.Meeting) -> dict[str, Any]:
    metadata = extract_tencent_meeting_metadata(meeting)
    records_args: dict[str, Any] = {"page_size": 10}
    if metadata.get("meeting_id"):
        records_args["meeting_id"] = metadata["meeting_id"]
    elif metadata.get("meeting_code"):
        records_args["meeting_code"] = metadata["meeting_code"]
    elif meeting.date:
        start = meeting.date.replace(hour=0, minute=0, second=0, microsecond=0)
        end = meeting.date.replace(hour=23, minute=59, second=59, microsecond=0)
        records_args["start_time"] = f"{start.isoformat()}+08:00"
        records_args["end_time"] = f"{end.isoformat()}+08:00"
    else:
        raise RuntimeError("会议卡片里没有腾讯会议号或会议时间，无法同步录制和转写")
    return records_args


def sync_tencent_meeting_minutes(db: Session, meeting: models.Meeting) -> models.Meeting:
    meeting.sync_status = "syncing"
    meeting.sync_error = ""
    db.commit()
    records_args = build_tencent_records_query(meeting)
    records_text = call_tencent_meeting_tool("get_records_list", records_args)
    record = parse_tencent_record_result(records_text)
    records = record.get("records") or [record]
    if not any(item.get("record_file_id") for item in records):
        raise TencentMinutesUnavailableError("暂未查询到腾讯会议录制文件，请确认会议已结束并已开启云录制/转写")

    transcript = ""
    selected_record = record
    for candidate in records:
        record_file_id = candidate.get("record_file_id")
        if not record_file_id:
            continue
        try:
            transcript = fetch_complete_tencent_transcript(record_file_id)
        except Exception:
            transcript = ""
        if transcript:
            selected_record = candidate
            break
    if not transcript:
        raise TencentMinutesUnavailableError("腾讯会议转写暂未生成，请稍后再同步")

    meeting.transcript = transcript
    link_parts = [
        part
        for part in (meeting.recording_url or "").splitlines()
        if part.strip()
        and not part.strip().startswith("meeting_record_id")
        and not part.strip().startswith("record_file_id")
        and not part.strip().startswith("录制下载")
    ]
    if selected_record.get("meeting_record_id"):
        link_parts.append(f"meeting_record_id：{selected_record['meeting_record_id']}")
        address_text = call_tencent_meeting_tool(
            "get_record_addresses",
            {"meeting_record_id": selected_record["meeting_record_id"]},
        )
        address = parse_tencent_record_address_result(address_text)
        if address["view_address"]:
            link_parts.append(f"录屏查看：{address['view_address']}")
            meeting.recording_view_url = address["view_address"]
        if address["trace"]:
            link_parts.append(f"录屏 X-Tc-Trace：{address['trace']}")
        if address["rpc_uuid"]:
            link_parts.append(f"录屏 rpcUuid：{address['rpc_uuid']}")
        meeting.sync_trace_json = json.dumps(
            {"X-Tc-Trace": address["trace"], "rpcUuid": address["rpc_uuid"]},
            ensure_ascii=False,
        )
    record_file_id = selected_record.get("record_file_id", "")
    if record_file_id:
        link_parts.append(f"record_file_id：{record_file_id}")
        meeting.record_file_id = record_file_id
    if record.get("download_url"):
        link_parts.append(f"录制下载：{record['download_url']}")
    meeting.recording_url = "\n".join(part for part in link_parts if part).strip()
    meeting.status = "transcribed"
    meeting.sync_status = "synced"
    meeting.sync_error = ""
    meeting.last_synced_at = datetime.now()
    project = db.get(models.Project, meeting.project_id)
    if project:
        deposit_meeting_transcript_to_knowledge(db, project, meeting)
    db.commit()
    db.refresh(meeting)
    return meeting

TEAM_MEMBERS = [
    ("项目经理", "项目负责人", ["业主沟通", "会议推进", "任务协调"], 45),
    ("主创负责人", "主创设计师", ["概念判断", "方案把控", "评审"], 55),
    ("总图负责人", "总图负责人", ["强排", "指标复核", "退界"], 50),
    ("立面负责人", "立面负责人", ["风格研究", "材料策略", "比例控制"], 35),
    ("汇报负责人", "汇报助理", ["PPT结构", "文本组织", "成果整合"], 40),
]


def ensure_digital_employees(db: Session) -> None:
    existing_names = set(db.scalars(select(models.DigitalEmployee.name)))
    target_names = {item[0] for item in DIGITAL_EMPLOYEES}
    if target_names.issubset(existing_names):
        return
    db.execute(delete(models.DigitalEmployee))
    for name, role, skills, avatar in DIGITAL_EMPLOYEES:
        db.add(
            models.DigitalEmployee(
                name=name,
                role=role,
                skills=json.dumps(skills, ensure_ascii=False),
                avatar=avatar,
                status="available",
                workload=20,
            )
        )
    db.commit()


def ensure_team_members(db: Session) -> None:
    existing_names = set(db.scalars(select(models.TeamMember.name)))
    if existing_names:
        return
    for name, role, skills, workload in TEAM_MEMBERS:
        db.add(
            models.TeamMember(
                name=name,
                role=role,
                skills=json.dumps(skills, ensure_ascii=False),
                status="available",
                workload=workload,
            )
        )
    db.commit()


def parse_document(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    try:
        if ext in {".png", ".jpg", ".jpeg"}:
            return "", "saved_no_ocr"
        if ext in {".txt", ".md"}:
            return path.read_text(encoding="utf-8", errors="ignore"), "parsed"
        if ext == ".pdf":
            reader = PdfReader(str(path))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
            return text, "parsed"
        if ext == ".docx":
            doc = Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
            return text, "parsed"
        if ext == ".xlsx":
            wb = load_workbook(str(path), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    values = [str(value) for value in row if value is not None]
                    if values:
                        rows.append(" | ".join(values))
            return "\n".join(rows), "parsed"
        if ext == ".pptx":
            try:
                from pptx import Presentation
            except Exception as exc:
                return f"解析失败：缺少 python-pptx（{exc}）", "failed"
            prs = Presentation(str(path))
            slides = []
            for idx, slide in enumerate(prs.slides, start=1):
                slides.append(f"# Slide {idx}")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        slides.append(text)
            return "\n".join(slides), "parsed"
        return "", "unsupported"
    except Exception as exc:
        return f"解析失败：{exc}", "failed"


def project_upload_dir(project_id: str) -> Path:
    path = settings.upload_root_path / "projects" / project_id
    path.mkdir(parents=True, exist_ok=True)
    return path


def project_managed_dir(project_id: str, upload_root: Optional[Path] = None) -> Path:
    root = upload_root or settings.upload_root_path
    return root / "projects" / project_id


def inbox_upload_dir() -> Path:
    path = settings.upload_root_path / "inbox"
    path.mkdir(parents=True, exist_ok=True)
    return path


def delete_project_library(db: Session, project_id: str, upload_root: Optional[Path] = None) -> dict[str, Any]:
    project = db.get(models.Project, project_id)
    if not project:
        return {"deleted": False, "deleted_project_id": project_id, "deleted_files": 0}
    managed_dir = project_managed_dir(project_id, upload_root)
    deleted_files = len(project.files)
    inbox_items = list(db.scalars(select(models.InboxItem).where(models.InboxItem.project_id == project_id)))
    for item in inbox_items:
        item.project_id = ""
        item.archive_path = ""
        item.status = "待确认"
        item.archive_group = "待确认"
        item.needs_review = True
    db.delete(project)
    db.commit()
    if managed_dir.exists() and managed_dir.is_dir():
        shutil.rmtree(managed_dir, ignore_errors=True)
    return {"deleted": True, "deleted_project_id": project_id, "deleted_files": deleted_files}


def _unique_path(directory: Path, filename: str) -> Path:
    directory.mkdir(parents=True, exist_ok=True)
    safe_name = Path(filename).name or "uploaded-file"
    candidate = directory / safe_name
    if not candidate.exists():
        return candidate
    stem = candidate.stem
    suffix = candidate.suffix
    for index in range(2, 1000):
        next_candidate = directory / f"{stem}_{index}{suffix}"
        if not next_candidate.exists():
            return next_candidate
    return directory / f"{stem}_{uuid4().hex[:8]}{suffix}"


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def create_inbox_item_from_path(db: Session, source_path: Path, target_dir: Optional[Path] = None, source_label: str = "本地路径") -> models.InboxItem:
    source_path = source_path.expanduser().resolve()
    if not source_path.exists() or not source_path.is_file():
        raise FileNotFoundError(f"文件不存在：{source_path}")
    ext = source_path.suffix.lower()
    if ext not in SUPPORTED_INBOX_EXTS:
        raise ValueError(f"不支持的文件类型：{source_path.name}")
    target = _unique_path(target_dir or inbox_upload_dir(), source_path.name)
    shutil.copy2(source_path, target)
    item = models.InboxItem(
        original_filename=source_path.name,
        suggested_filename=source_path.name,
        final_filename="",
        source_path=str(source_path),
        temp_path=str(target),
        source_label=source_label,
        status="待确认",
        needs_review=True,
        file_hash=file_sha256(target),
        archive_group="待确认",
    )
    db.add(item)
    db.commit()
    db.refresh(item)
    return item


def run_inbox_scan_with_progress(
    db: Session,
    root: Path,
    source_label: str,
    days: int = 0,
    progress: Optional[dict[str, Any]] = None,
) -> dict[str, Any]:
    root = root.expanduser()
    if not root.exists() or not root.is_dir():
        raise FileNotFoundError(f"目录不存在：{root}")
    cutoff = datetime.now().timestamp() - (days * 24 * 60 * 60) if days and days > 0 else 0
    if progress is not None:
        progress.update(
            {
                "status": "running",
                "step": "枚举文件",
                "total_candidates": 0,
                "processed": 0,
                "imported_files": 0,
                "unsupported_files": 0,
                "old_files": 0,
                "failed_files": 0,
                "current_file": "",
            }
        )

    candidates: list[Path] = []
    unsupported = 0
    old_files = 0
    for path in sorted(root.rglob("*"), key=lambda item: item.stat().st_mtime if item.exists() else 0, reverse=True):
        if len(candidates) >= 200:
            break
        if not path.is_file():
            continue
        if path.suffix.lower() not in SUPPORTED_INBOX_EXTS:
            unsupported += 1
            continue
        if path.stat().st_mtime < cutoff:
            old_files += 1
            continue
        candidates.append(path)

    if progress is not None:
        progress.update(
            {
                "step": "复制并识别",
                "total_candidates": len(candidates),
                "unsupported_files": unsupported,
                "old_files": old_files,
            }
        )

    imported: list[models.InboxItem] = []
    failed: list[str] = []
    for index, path in enumerate(candidates, start=1):
        if progress is not None:
            progress.update({"processed": index, "current_file": path.name, "step": "复制到收件箱"})
        try:
            item = create_inbox_item_from_path(db, path, inbox_upload_dir(), source_label)
            if progress is not None:
                progress.update({"step": "识别项目和资料类型"})
            imported.append(classify_inbox_item(db, item))
            if progress is not None:
                progress["imported_files"] = len(imported)
        except Exception as exc:
            failed.append(f"{path.name}：{exc}")
            if progress is not None:
                progress["failed_files"] = len(failed)

    if not imported:
        detail = f"未导入文件。目录中不支持格式 {unsupported} 个"
        if old_files:
            detail += f"，超过时间范围 {old_files} 个"
        if failed:
            detail += "，失败：" + "；".join(failed[:5])
        raise ValueError(detail)

    if progress is not None:
        progress.update({"step": "生成整体建议", "current_file": ""})
    advice = build_inbox_batch_advice(db, [item.id for item in imported])
    result = {
        "imported_count": len(imported),
        "items": imported,
        "unsupported_files": unsupported,
        "old_files": old_files,
        "failed_files": failed,
        "batch_advice": advice,
    }
    if progress is not None:
        progress.update(
            {
                "status": "succeeded",
                "step": "完成",
                "processed": len(candidates),
                "current_file": "",
                "result": {
                    "imported_count": len(imported),
                    "unsupported_files": unsupported,
                    "old_files": old_files,
                    "failed_count": len(failed),
                    "batch_advice": advice,
                },
            }
        )
    return result


def _read_small_context(path: Path, limit: int = 4000) -> str:
    if path.suffix.lower() not in SUPPORTED_PROJECT_EXTS:
        return ""
    text, status = parse_document(path)
    if status != "parsed":
        return ""
    return text[:limit]


def _pick_material_type(text: str, ext: str) -> tuple[str, list[str]]:
    rules = [
        ("会议资料", ["会议纪要", "参会人", "议题", "会议录音", "转写", "启动会"]),
        ("技术条件", ["日照", "退界", "容积率", "消防", "规划条件", "人防", "绿建", "管线"]),
        ("参考案例", ["竞品", "案例", "参考", "市场定位"]),
        ("审核反馈", ["甲方意见", "修改意见", "审核", "反馈", "批注", "问题清单"]),
        ("交付成果", ["汇报", "提案", "投标", "正式提交", "盖章版", "交付"]),
        ("项目基础资料", ["设计任务书", "红线", "面积表", "合同", "甲方需求", "用地"]),
    ]
    if ext in {".zip", ".rar", ".7z"}:
        return "压缩包与杂项", ["压缩包待确认"]
    if ext in {".dwg", ".skp", ".rvt"}:
        return "设计源文件", ["设计源文件仅登记资产"]
    if ext in {".png", ".jpg", ".jpeg", ".webp"}:
        return "设计过程", ["图片资料默认先做项目资产"]
    for material_type, keywords in rules:
        hits = [word for word in keywords if word in text]
        if hits:
            return material_type, hits
    if ext in {".pptx", ".pdf"}:
        return "设计过程", ["按文件类型判断"]
    return "项目基础资料", ["默认归入基础资料"]


PROJECT_TOPIC_WORDS = {
    "项目汇报",
    "汇报",
    "方案设计",
    "概念方案设计",
    "概念方案",
    "方案",
    "定位报告",
    "总图及户型",
    "总图",
    "户型",
    "附件",
    "配套",
    "竞品",
    "集",
    "最终版",
    "新定案",
}


def _strip_project_topic_words(text: str) -> str:
    value = re.sub(r"\.[^.]+$", "", text)
    value = re.sub(r"20\d{2}[-年.]?\d{1,2}[-月.]?\d{0,2}日?", "", value)
    value = re.sub(r"26\d{4}", "", value)
    value = re.sub(r"0115|0103|0114|0118|V\d+|v\d+", "", value)
    value = re.sub(r"附件\s*\d*[：:、-]?", "", value)
    value = re.sub(r"[（(][^）)]*[）)]", "", value)
    for word in sorted(PROJECT_TOPIC_WORDS, key=len, reverse=True):
        value = value.replace(word, "")
    value = re.sub(r"[_\-\s]+", "", value)
    value = re.sub(r"[：:、，,。]+", "", value)
    return value.strip()


def _canonical_project_name_from_text(filename: str, context: str = "") -> str:
    haystack = f"{filename}\n{context}"
    if "振三街" in haystack:
        if "石家庄" in haystack:
            return "石家庄市振三街地块项目"
        return "振三街项目"
    patterns = [
        r"([\u4e00-\u9fa5]{2,12}(?:市|区|县)?[\u4e00-\u9fa5A-Za-z0-9]{1,18}(?:地块|项目))",
        r"([\u4e00-\u9fa5A-Za-z0-9]{2,24}(?:项目|地块))",
    ]
    for pattern in patterns:
        match = re.search(pattern, haystack)
        if match:
            candidate = _strip_project_topic_words(match.group(1))
            if candidate:
                if candidate.endswith("地块"):
                    return f"{candidate}项目"
                return candidate[:32]
    stem = Path(filename).stem
    cleaned = _strip_project_topic_words(stem)
    parts = [part for part in re.split(r"[_\-\s]+", cleaned) if part]
    stop_words = {"启动会", "会议纪要", "规划条件", "日照退界", "日照分析", "设计任务书"}
    useful = [part for part in parts if part not in stop_words]
    return (useful[0] if useful else cleaned or stem)[:32]


def _extract_project_hint(filename: str, context: str = "") -> str:
    return _canonical_project_name_from_text(filename, context)


def _extract_date_token(text: str) -> str:
    match = re.search(r"(20\d{2})[-年.]?(\d{1,2})[-月.]?(\d{1,2})?", text)
    if not match:
        return datetime.now().strftime("%Y%m%d")
    year, month, day = match.group(1), match.group(2), match.group(3) or "01"
    return f"{year}{int(month):02d}{int(day):02d}"


def _safe_filename_part(text: str, fallback: str = "资料") -> str:
    value = re.sub(r"[\\/:*?\"<>|\s]+", "", text or "")
    return value[:32] or fallback


KNOWLEDGE_RECOMMENDED_TYPES = {"技术条件", "会议资料", "审核反馈", "项目基础资料", "参考案例", "复盘沉淀"}
KNOWLEDGE_SKIPPED_REASONS = {
    "设计源文件": "设计源文件默认只做项目资产记录。",
    "压缩包与杂项": "压缩包与杂项默认只做项目资产记录。",
    "设计过程": "设计过程稿默认先归档到项目，确认为可复用成果后再入知识库。",
}


def _knowledge_recommendation(material_type: str, filename: str = "", text: str = "") -> tuple[bool, str]:
    haystack = f"{filename}\n{text}"
    if material_type in {"设计源文件", "压缩包与杂项", "设计过程"}:
        return False, KNOWLEDGE_SKIPPED_REASONS.get(material_type, "该资料默认只归档到项目。")
    project_specific_words = ["项目汇报", "方案设计", "概念方案", "总图", "户型", "阶段性方案", "汇报稿"]
    reusable_rules = [
        ("规划条件", "包含规划条件，具备跨项目复用价值，可供后续项目参考。"),
        ("退界", "包含退界要求，具备跨项目复用价值，可供后续项目参考。"),
        ("日照", "包含日照要求，具备跨项目复用价值，可供后续项目参考。"),
        ("消防", "包含消防条件，具备跨项目复用价值，可供后续项目参考。"),
        ("人防", "包含人防条件，具备跨项目复用价值，可供后续项目参考。"),
        ("绿建", "包含绿建要求，具备跨项目复用价值，可供后续项目参考。"),
        ("建筑面积计算规则", "包含地方规则或面积计算规则，建议进入知识库。"),
        ("会议结论", "包含会议结论，建议进入知识库。"),
        ("审核意见", "包含审核意见，建议进入知识库。"),
        ("甲方反馈", "包含甲方反馈，建议进入知识库。"),
        ("复盘", "包含复盘方法，建议进入知识库。"),
        ("竞品", "包含竞品或市场定位结论，建议进入知识库。"),
        ("市场定位", "包含竞品或市场定位结论，建议进入知识库。"),
        ("定位报告", "包含竞品或市场定位结论，建议进入知识库。"),
    ]
    if any(word in filename for word in project_specific_words):
        if any(word in haystack for word in ["规划条件", "建筑面积计算规则", "审核意见", "甲方反馈", "会议结论", "复盘", "竞品", "市场定位", "定位报告"]):
            for word, reason in reusable_rules:
                if word in haystack:
                    return True, f"建议入库：{reason}"
        return False, "仅项目归档：这是本项目阶段性方案或汇报资料，复用价值有限。"
    for word, reason in reusable_rules:
        if word in haystack:
            return True, f"建议入库：{reason}"
    if material_type in {"会议资料", "审核反馈", "参考案例"}:
        return True, f"建议入库：{material_type}通常具备复用价值。"
    return False, "仅项目归档：未识别到明确可复用的技术条件、反馈、复盘或规则内容。"


def _same_file_hash(path: str, expected_hash: str) -> bool:
    if not path or not expected_hash:
        return False
    candidate = Path(path)
    if not candidate.exists() or not candidate.is_file():
        return False
    try:
        return file_sha256(candidate) == expected_hash
    except Exception:
        return False


def _detect_duplicate(db: Session, item: models.InboxItem) -> tuple[str, str, str]:
    file_hash = item.file_hash or (file_sha256(Path(item.temp_path)) if item.temp_path else "")
    if not file_hash:
        return "", "", ""
    for project_file in db.scalars(select(models.ProjectFile)):
        if _same_file_hash(project_file.filepath, file_hash):
            return "project", project_file.id, ""
    for knowledge_file in db.scalars(select(models.KnowledgeFile)):
        if _same_file_hash(knowledge_file.filepath, file_hash):
            return "knowledge", "", knowledge_file.id
    return "", "", ""


def recommend_inbox_item(db: Session, item: models.InboxItem) -> models.InboxItem:
    if item.temp_path and Path(item.temp_path).exists():
        item.file_hash = item.file_hash or file_sha256(Path(item.temp_path))
    duplicate_scope, duplicate_project_file_id, duplicate_knowledge_file_id = _detect_duplicate(db, item)
    item.duplicate_scope = duplicate_scope
    item.duplicate_project_file_id = duplicate_project_file_id
    item.duplicate_knowledge_file_id = duplicate_knowledge_file_id
    if duplicate_scope:
        item.status = "重复文件"
        item.archive_group = "重复文件"
        item.recommended_action = "重复跳过"
        item.suggest_knowledge = False
        item.recommend_knowledge_reason = "项目库或知识库已存在相同文件，默认跳过。"
        item.needs_review = True
        db.commit()
        db.refresh(item)
        return item

    suggest_knowledge, reason = _knowledge_recommendation(item.material_type, item.original_filename, item.summary)
    item.suggest_knowledge = suggest_knowledge
    item.recommend_knowledge_reason = reason
    if item.project_id and item.confidence >= 0.85:
        item.archive_group = "可直接确认"
        item.recommended_action = "入知识库" if suggest_knowledge else "仅项目归档"
        item.status = "待确认"
        item.needs_review = False
    elif not item.project_id:
        item.archive_group = "未归属项目"
        item.recommended_action = "创建新项目"
        item.status = "未归属项目"
        item.needs_review = True
    else:
        item.archive_group = "需审核"
        item.recommended_action = "需人工确认"
        item.status = "需审核"
        item.needs_review = True
    db.commit()
    db.refresh(item)
    return item


def classify_inbox_item(db: Session, item: models.InboxItem) -> models.InboxItem:
    path = Path(item.temp_path)
    ext = path.suffix.lower()
    context = _read_small_context(path)
    haystack = f"{item.original_filename}\n{context}"
    projects = list(db.scalars(select(models.Project).order_by(models.Project.updated_at.desc())))
    matched_project = None
    evidence: list[str] = []
    confidence = 0.35
    for project in projects:
        if project.name and project.name in haystack:
            matched_project = project
            evidence.append(f"文件名或内容包含项目名：{project.name}")
            confidence = 0.92
            break
        if project.city and project.city in haystack and project.phase and project.phase in haystack:
            matched_project = project
            evidence.append(f"同时匹配城市和阶段：{project.city}/{project.phase}")
            confidence = 0.72
            break

    material_type, keyword_hits = _pick_material_type(haystack, ext)
    evidence.extend([f"关键词：{word}" for word in keyword_hits[:4]])
    project_name = matched_project.name if matched_project else _extract_project_hint(item.original_filename, context)
    phase = matched_project.phase if matched_project else ("强排" if "强排" in haystack else "待确认")
    topic = keyword_hits[0] if keyword_hits else Path(item.original_filename).stem[:12]
    date_token = _extract_date_token(item.original_filename)
    suggested_filename = "_".join(
        [
            _safe_filename_part(project_name, "未归属项目"),
            _safe_filename_part(material_type),
            _safe_filename_part(topic),
            _safe_filename_part(phase),
            date_token,
            "V1",
        ]
    ) + ext

    item.project_id = matched_project.id if matched_project else ""
    item.suggested_project_name = project_name
    item.suggested_city = matched_project.city if matched_project else ""
    item.suggested_project_type = matched_project.project_type if matched_project else ""
    item.suggested_phase = phase
    item.material_type = material_type
    item.summary = (context.strip().replace("\n", " ")[:180] if context else f"{material_type}文件，等待人工确认归档。")
    item.keywords = json.dumps(keyword_hits, ensure_ascii=False)
    item.evidence = "；".join(evidence)
    item.confidence = confidence
    item.status = "待确认" if matched_project else "未归属项目"
    item.needs_review = confidence < 0.85 or not matched_project
    suggest_knowledge, knowledge_reason = _knowledge_recommendation(material_type, item.original_filename, context)
    item.suggest_knowledge = suggest_knowledge
    item.recommend_knowledge_reason = knowledge_reason
    item.suggest_todo = material_type in {"会议资料", "审核反馈"}
    item.suggested_filename = suggested_filename
    if path.exists():
        item.file_hash = item.file_hash or file_sha256(path)
    db.commit()
    db.refresh(item)
    return recommend_inbox_item(db, item)


def delete_inbox_items(db: Session, item_ids: list[str]) -> int:
    items = list(db.scalars(select(models.InboxItem).where(models.InboxItem.id.in_(item_ids))))
    for item in items:
        path = Path(item.temp_path) if item.temp_path else None
        if path and path.exists() and path.is_file():
            try:
                path.unlink()
            except OSError:
                pass
        db.delete(item)
    db.commit()
    return len(items)


def apply_inbox_items(
    db: Session,
    item_ids: list[str],
    project_id: str = "",
    project_payload: Optional[Any] = None,
    final_filename_by_id: Optional[dict[str, str]] = None,
    material_type_by_id: Optional[dict[str, str]] = None,
    enter_knowledge: bool = False,
    enter_knowledge_by_id: Optional[dict[str, bool]] = None,
    archive_root: Optional[Path] = None,
) -> dict[str, Any]:
    items = list(db.scalars(select(models.InboxItem).where(models.InboxItem.id.in_(item_ids))))
    if not items:
        raise ValueError("没有找到可归档的收件箱文件")
    project = db.get(models.Project, project_id) if project_id else None
    if not project:
        if project_payload is None:
            first = items[0]
            project_payload = type("ProjectPayload", (), {
                "name": first.suggested_project_name or "未命名项目",
                "city": first.suggested_city,
                "project_type": first.suggested_project_type,
                "phase": first.suggested_phase,
                "description": "由文件收件箱创建",
                "status": "active",
            })()
        if hasattr(project_payload, "model_dump"):
            data = project_payload.model_dump()
        elif isinstance(project_payload, dict):
            data = project_payload
        else:
            data = {
                "name": getattr(project_payload, "name", "未命名项目"),
                "city": getattr(project_payload, "city", ""),
                "project_type": getattr(project_payload, "project_type", ""),
                "phase": getattr(project_payload, "phase", ""),
                "description": getattr(project_payload, "description", ""),
                "status": getattr(project_payload, "status", "active"),
            }
        project = models.Project(**data)
        db.add(project)
        db.commit()
        db.refresh(project)

    target_dir = archive_root or project_upload_dir(project.id)
    target_dir.mkdir(parents=True, exist_ok=True)
    final_filename_by_id = final_filename_by_id or {}
    material_type_by_id = material_type_by_id or {}
    files: list[models.ProjectFile] = []
    for item in items:
        final_name = Path(final_filename_by_id.get(item.id) or item.suggested_filename or item.original_filename).name
        if not Path(final_name).suffix:
            final_name += Path(item.original_filename).suffix
        target = _unique_path(target_dir, final_name)
        shutil.copy2(Path(item.temp_path), target)
        text, parse_status = parse_document(target) if target.suffix.lower() in SUPPORTED_PROJECT_EXTS else ("", "saved_no_parse")
        record = models.ProjectFile(
            project_id=project.id,
            filename=target.name,
            filepath=str(target),
            filetype=target.suffix.lower().lstrip("."),
            filesize=target.stat().st_size,
            parsed_text=text,
            parse_status=parse_status,
            analysis_status="pending",
        )
        db.add(record)
        files.append(record)
        item.project_id = project.id
        item.final_filename = target.name
        item.archive_path = str(target)
        item.material_type = material_type_by_id.get(item.id, item.material_type)
        item.status = "已归档"
        item.needs_review = False
        should_enter_knowledge = enter_knowledge or bool((enter_knowledge_by_id or {}).get(item.id))
        if should_enter_knowledge:
            index_knowledge_file(db, target, display_path=f"projects/{project.id}/{target.name}")
            item.status = "已进入知识库"
        item.archive_group = item.status
    db.commit()
    for file in files:
        db.refresh(file)
    for item in items:
        db.refresh(item)
    db.refresh(project)
    return {"project": project, "files": files, "items": items}


def apply_inbox_recommendations(
    db: Session,
    item_ids: list[str],
    force_duplicate_ids: Optional[list[str]] = None,
    archive_root: Optional[Path] = None,
) -> dict[str, Any]:
    force_duplicate_ids = force_duplicate_ids or []
    items = list(db.scalars(select(models.InboxItem).where(models.InboxItem.id.in_(item_ids))))
    files: list[models.ProjectFile] = []
    changed_items: list[models.InboxItem] = []
    skipped_count = 0
    created_project_count = 0
    actionable_groups: dict[tuple[str, str, str], list[models.InboxItem]] = defaultdict(list)
    batch_keys = _batch_project_group_keys(items)
    for item in items:
        item = recommend_inbox_item(db, item)
        if item.archive_group == "重复文件" and item.id not in force_duplicate_ids:
            skipped_count += 1
            changed_items.append(item)
            continue
        if item.recommended_action == "需人工确认":
            skipped_count += 1
            changed_items.append(item)
            continue
        actionable_groups[batch_keys.get(item.id, _item_project_group_key(item))].append(item)

    for group_key, group_items in actionable_groups.items():
        kind, _, project_name = group_key
        project_id = group_items[0].project_id if kind == "existing" else ""
        payload = None
        if not project_id:
            payload = type("ProjectPayload", (), {
                "name": project_name or "未命名项目",
                "city": group_items[0].suggested_city,
                "project_type": group_items[0].suggested_project_type or "住宅",
                "phase": group_items[0].suggested_phase or "待确认",
                "description": f"由文件收件箱根据 {len(group_items)} 个资料创建",
                "status": "active",
            })()
        result = apply_inbox_items(
            db,
            [item.id for item in group_items],
            project_id=project_id,
            project_payload=payload,
            enter_knowledge_by_id={item.id: item.suggest_knowledge for item in group_items},
            archive_root=archive_root,
        )
        files.extend(result["files"])
        changed_items.extend(result["items"])
        if not project_id and result["project"].id:
            created_project_count += 1
    return {
        "files": files,
        "items": changed_items,
        "skipped_count": skipped_count,
        "created_project_count": created_project_count,
    }


def _item_project_group_key(item: models.InboxItem) -> tuple[str, str, str]:
    if item.project_id:
        return ("existing", item.project_id, item.suggested_project_name or "已有项目")
    canonical = _canonical_project_name_from_text(item.original_filename, item.summary)
    return ("new", canonical or item.suggested_project_name or "未命名项目", canonical or item.suggested_project_name or "未命名项目")


def _project_signature(name: str) -> str:
    if "振三街" in name:
        return "振三街"
    value = re.sub(r"(?:项目|地块|强排|方案|设计|总图|户型|规划条件|启动会|会议纪要)", "", name)
    value = re.sub(r"(?:杭州|萧山|杭州市|萧山区|浙江|浙江省|市|区|县)", "", value)
    if len(value) >= 2:
        return value
    value = re.sub(r"(?:^.*市|^.*区|^.*县)", "", name)
    value = re.sub(r"(项目|地块)$", "", value)
    return value or name


def _batch_project_group_keys(items: list[models.InboxItem]) -> dict[str, tuple[str, str, str]]:
    raw_keys = {item.id: _item_project_group_key(item) for item in items}
    by_signature: dict[str, list[tuple[str, str, str]]] = defaultdict(list)
    for key in raw_keys.values():
        kind, project_key, project_name = key
        if kind == "new":
            by_signature[_project_signature(project_name)].append(key)
    preferred_by_signature: dict[str, tuple[str, str, str]] = {}
    for signature, keys in by_signature.items():
        preferred = sorted(set(keys), key=lambda key: (len(key[2]), "市" in key[2], "地块" in key[2]), reverse=True)[0]
        preferred_by_signature[signature] = preferred
    return {
        item_id: preferred_by_signature.get(_project_signature(key[2]), key) if key[0] == "new" else key
        for item_id, key in raw_keys.items()
    }


def _build_inbox_advice_markdown(
    total_files: int,
    action_counts: dict[str, int],
    project_groups: list[dict[str, Any]],
    knowledge_candidates: list[dict[str, str]],
    duplicates: list[dict[str, str]],
    needs_review: list[dict[str, str]],
) -> str:
    lines = [
        f"本次收件箱共有 {total_files} 个文件。",
        f"建议直接归档 {action_counts['归档文件']} 个，其中需要创建新项目 {action_counts['创建项目']} 个，建议进入知识库 {action_counts['入知识库']} 个。",
    ]
    if duplicates:
        lines.append(f"发现 {len(duplicates)} 个重复文件，默认跳过，避免重复进入项目库或知识库。")
    if needs_review:
        lines.append(f"还有 {len(needs_review)} 个文件需要人工看一眼，主要是项目归属或资料类型不够确定。")
    if project_groups:
        lines.append("按项目归档建议：")
        for group in project_groups[:8]:
            action = "建议创建" if group.get("kind") == "new" else "归入已有"
            lines.append(f"- {action}：{group['project_name']}，{group['file_count']} 个文件，包含 {group['material_summary']}。")
    if knowledge_candidates:
        lines.append(f"知识库建议：优先收录 {len(knowledge_candidates)} 个可复用资料，例如技术条件、会议资料、审核反馈、项目基础资料。")
    return "\n".join(lines)


def build_inbox_batch_advice(db: Session, item_ids: list[str]) -> dict[str, Any]:
    query = select(models.InboxItem).order_by(models.InboxItem.created_at.desc())
    if item_ids:
        query = query.where(models.InboxItem.id.in_(item_ids))
    items = list(db.scalars(query))

    recommended_item_ids: list[str] = []
    knowledge_candidates: list[dict[str, str]] = []
    duplicates: list[dict[str, str]] = []
    needs_review: list[dict[str, str]] = []
    grouped_items: dict[tuple[str, str, str], list[models.InboxItem]] = defaultdict(list)
    create_project_keys: set[tuple[str, str, str]] = set()
    batch_keys = _batch_project_group_keys(items)

    for item in items:
        item = recommend_inbox_item(db, item)
        if item.archive_group == "重复文件" or item.recommended_action == "重复跳过":
            duplicates.append(
                {
                    "id": item.id,
                    "filename": item.original_filename,
                    "scope": item.duplicate_scope or "unknown",
                    "reason": item.recommend_knowledge_reason or "已存在相同文件",
                }
            )
            continue
        if item.recommended_action == "需人工确认" or item.archive_group == "需审核":
            needs_review.append(
                {
                    "id": item.id,
                    "filename": item.original_filename,
                    "reason": item.evidence or "项目归属或资料类型不够确定",
                }
            )
            continue

        recommended_item_ids.append(item.id)
        group_key = batch_keys.get(item.id, _item_project_group_key(item))
        grouped_items[group_key].append(item)
        if group_key[0] == "new":
            create_project_keys.add(group_key)
        if item.suggest_knowledge:
            knowledge_candidates.append(
                {
                    "id": item.id,
                    "filename": item.original_filename,
                    "material_type": item.material_type,
                    "reason": item.recommend_knowledge_reason,
                }
            )

    project_groups: list[dict[str, Any]] = []
    for (kind, project_key, project_name), group_items in grouped_items.items():
        material_counts = Counter(item.material_type or "未分类" for item in group_items)
        project_groups.append(
            {
                "kind": kind,
                "project_id": project_key if kind == "existing" else "",
                "project_name": project_name,
                "file_count": len(group_items),
                "knowledge_count": sum(1 for item in group_items if item.suggest_knowledge),
                "item_ids": [item.id for item in group_items],
                "aliases": sorted({item.suggested_project_name for item in group_items if item.suggested_project_name})[:5],
                "material_summary": "、".join(f"{name}{count}个" for name, count in material_counts.most_common(4)),
            }
        )
    project_groups.sort(key=lambda group: group["file_count"], reverse=True)

    action_counts = {
        "归档文件": len(recommended_item_ids),
        "创建项目": len(create_project_keys),
        "入知识库": len(knowledge_candidates),
        "跳过重复": len(duplicates),
        "需审核": len(needs_review),
    }
    markdown = _build_inbox_advice_markdown(
        len(items),
        action_counts,
        project_groups,
        knowledge_candidates,
        duplicates,
        needs_review,
    )
    return {
        "total_files": len(items),
        "recommended_item_ids": recommended_item_ids,
        "action_counts": action_counts,
        "project_groups": project_groups,
        "knowledge_candidates": knowledge_candidates,
        "duplicates": duplicates,
        "needs_review": needs_review,
        "markdown": markdown,
        "mode": "rule",
    }


async def build_inbox_batch_advice_with_ai(db: Session, item_ids: list[str]) -> dict[str, Any]:
    advice = build_inbox_batch_advice(db, item_ids)
    if settings.mock_mode or not settings.deepseek_api_key:
        return advice

    query = select(models.InboxItem).order_by(models.InboxItem.created_at.desc())
    if item_ids:
        query = query.where(models.InboxItem.id.in_(item_ids))
    items = list(db.scalars(query))
    file_rows = [
        {
            "文件名": item.original_filename,
            "建议项目": item.suggested_project_name,
            "资料类型": item.material_type,
            "推荐动作": item.recommended_action,
            "是否入知识库": item.suggest_knowledge,
            "重复状态": item.duplicate_scope,
            "摘要": (item.summary or "")[:160],
            "依据": (item.evidence or "")[:160],
        }
        for item in items
    ]
    prompt = f"""
你是 ROM-AI 的项目资料管理员。请根据下面收件箱文件清单，输出一份给项目经理看的整体归档建议。

要求：
1. 不要逐个文件长篇解释，要按项目和资料角色分组。
2. 说明哪些可以直接归档，哪些建议创建新项目，哪些建议进入知识库，哪些需要人工审核。
3. 给出本次归档的总体判断和风险提醒。
4. 不要改变系统已计算的动作，只解释和组织建议。

系统统计：
{json.dumps(advice["action_counts"], ensure_ascii=False)}

文件清单：
{json.dumps(file_rows, ensure_ascii=False)}
""".strip()
    try:
        ai_markdown = await call_deepseek_text(
            prompt,
            "你负责把项目文件收件箱的识别结果整理成清晰、可执行、非技术化的归档建议。",
        )
    except Exception:
        return advice
    if ai_markdown.strip():
        advice["markdown"] = ai_markdown.strip()
        advice["mode"] = "deepseek"
    return advice


def mock_analysis_payload(project: models.Project, files: list[models.ProjectFile]) -> dict[str, Any]:
    completeness = "低" if not files else "中"
    tasks = [
        {
            "task_name": "宋式展示区归家动线策略",
            "task_type": "策略分析",
            "priority": "高",
            "owner_role": "主创设计师",
            "estimated_days": 3,
            "dependencies": ["T001"],
            "risk_level": "中",
            "status": "todo",
            "output_requirement": "输出动线分析图、策略文本、汇报PPT结构。",
        },
        {
            "task_name": "社区礼序界面与门头尺度校准",
            "task_type": "立面研究",
            "priority": "高",
            "owner_role": "立面负责人",
            "estimated_days": 4,
            "dependencies": ["宋式展示区归家动线策略"],
            "risk_level": "中",
            "status": "todo",
            "output_requirement": "输出门头比例、材质建议和节点风险说明。",
        },
        {
            "task_name": "展示区景观节点与室内到访体验衔接",
            "task_type": "协同设计",
            "priority": "中",
            "owner_role": "景观接口",
            "estimated_days": 2,
            "dependencies": ["宋式展示区归家动线策略"],
            "risk_level": "低",
            "status": "todo",
            "output_requirement": "输出景观节点清单、室内接口条件和体验动线说明。",
        },
    ]
    timeline = [
        {
            "stage_name": "概念方案阶段",
            "start_day": 1,
            "end_day": 7,
            "milestone": "概念方案汇报",
            "dependencies": [],
            "risk_note": "甲方决策周期不确定。",
        },
        {
            "stage_name": "展示区深化阶段",
            "start_day": 8,
            "end_day": 14,
            "milestone": "展示区动线与界面确认",
            "dependencies": ["概念方案阶段"],
            "risk_note": "景观、室内、立面接口需要同步确认。",
        },
        {
            "stage_name": "汇报整合阶段",
            "start_day": 15,
            "end_day": 18,
            "milestone": "PPT与图像成果交付",
            "dependencies": ["展示区深化阶段"],
            "risk_note": "素材版本与甲方关注点可能发生变化。",
        },
    ]
    team_requirements = {
        "total_headcount": 5,
        "roles": [
            {"role": "项目负责人", "count": 1, "skills": ["项目管理", "甲方沟通"], "intensity": "全职"},
            {"role": "主创设计师", "count": 1, "skills": ["宋式风格", "展示区设计"], "intensity": "全职"},
            {"role": "立面负责人", "count": 1, "skills": ["比例控制", "材料策略"], "intensity": "阶段投入"},
            {"role": "景观接口", "count": 1, "skills": ["归家动线", "节点体验"], "intensity": "阶段投入"},
            {"role": "汇报助理", "count": 1, "skills": ["PPT结构", "图像提示词"], "intensity": "阶段投入"},
        ],
    }
    knowledge_refs = [
        {
            "source_file": "宋式社区方法论.md",
            "source_path": "Obsidian Vault/方法/宋式社区方法论.md",
            "chunk_id": "mock-song-community-001",
            "quote": "宋式展示区应强调归家动线的礼仪感，以门庭、院落、廊下空间形成连续体验。",
            "relevance_score": 0.92,
        },
        {
            "source_file": "展示区体验动线清单.md",
            "source_path": "Obsidian Vault/方法/展示区体验动线清单.md",
            "chunk_id": "mock-arrival-002",
            "quote": "首开展示区需要把车行到达、步行归家、接待转换和样板间参观整合为一条可讲述的路径。",
            "relevance_score": 0.86,
        },
    ]
    return {
        "mode": "mock",
        "project_basis": {
            "project_type": project.project_type or "待确认",
            "phase": project.phase or "待确认",
            "资料完整度": completeness,
            "缺失信息": ["规划条件细则", "成本目标", "甲方决策边界"] if not files else ["成本目标", "关键节点时间"],
            "综合风险等级": "中",
        },
        "design_difficulties": {
            "技术难点": ["需要尽快核对指标、总图边界、产品组合和消防/日照等基础约束。"],
            "协调难点": ["建筑、景观、室内和报规口径需要统一，避免汇报阶段反复返工。"],
            "规范难点": ["消防、日照、停车和无障碍需要在强排阶段提前校核。"],
            "甲方决策难点": ["需要明确产品档次、立面成本和展示区范围。"],
            "成本与落地难点": ["立面材料、景观节点和公区精装需要控制落地成本。"],
            "后续深化风险点": ["资料不完整会影响任务拆解、周期判断和团队配置。"],
        },
        "timeline_summary": {
            "总体推进周期": "约 28-42 天（Mock，需要人工校准）",
            "概念方案阶段": "5-7 天",
            "强排 / 总图阶段": "7-10 天",
            "户型与产品阶段": "5-8 天",
            "立面与风格阶段": "7-10 天",
            "景观 / 室内 / 精装协同阶段": "5-8 天",
            "报规或汇报成果阶段": "3-5 天",
            "关键路径": ["规划条件确认", "强排稳定", "立面方向锁定", "汇报成果整合"],
            "里程碑节点": ["概念方向会", "强排评审", "立面评审", "最终汇报"],
            "可能延误点": ["资料缺失", "甲方反复", "多专业接口不同步"],
        },
        "timeline": timeline,
        "staffing": {
            "项目负责人 / PM": "1人",
            "主创设计师": "1人",
            "总图负责人": "1人",
            "户型负责人": "1人",
            "立面负责人": "1人",
            "景观接口": "0.5人",
            "室内接口": "0.5人",
            "后期 / 报规接口": "0.5人",
            "AI辅助人员": "1人",
            "建议人数规模": "5-7人等效投入",
            "每类人员技能要求": ["强排经验", "产品定位", "立面落地", "汇报组织"],
            "各阶段人员投入强度": "前期 PM/总图高，汇报前主创/立面/汇报助理高。",
        },
        "team_requirements": team_requirements,
        "knowledge_refs": knowledge_refs,
        "tasks": tasks,
        "next_actions": [
            "补齐规划条件、红线、指标和甲方任务书。",
            "建立项目资料目录并标注版本。",
            "先做强排风险清单，再进入立面风格推演。",
            "明确展示区、首开区和汇报成果范围。",
            "把关键问题提交甲方形成一次决策会。",
        ],
    }


def analysis_to_markdown(payload: dict[str, Any]) -> str:
    lines = [f"# 项目分析报告（{payload.get('mode', 'mock')}）"]
    for section, value in payload.items():
        if section == "mode":
            continue
        lines.append(f"\n## {section}")
        if isinstance(value, dict):
            for key, item in value.items():
                rendered = json.dumps(item, ensure_ascii=False) if not isinstance(item, str) else item
                lines.append(f"- **{key}**：{rendered}")
        elif isinstance(value, list):
            for item in value:
                rendered = json.dumps(item, ensure_ascii=False) if isinstance(item, dict) else item
                lines.append(f"- {rendered}")
        else:
            lines.append(str(value))
    return "\n".join(lines)


def normalize_analysis_payload(payload: dict[str, Any], fallback: dict[str, Any]) -> dict[str, Any]:
    merged = {**fallback, **payload}
    for key in ("tasks", "timeline", "knowledge_refs"):
        if not isinstance(merged.get(key), list) or not merged[key]:
            merged[key] = fallback.get(key, [])
    if not isinstance(merged.get("team_requirements"), dict):
        merged["team_requirements"] = fallback.get("team_requirements", {"total_headcount": 0, "roles": []})
    return merged


async def call_deepseek_json(prompt: str) -> dict[str, Any]:
    if settings.mock_mode:
        return {}
    headers = {"Authorization": f"Bearer {settings.deepseek_api_key}", "Content-Type": "application/json"}
    body = {
        "model": settings.deepseek_model,
        "temperature": 0.2,
        "messages": [
            {
                "role": "system",
                "content": "你是建筑设计项目分析助手。只输出 JSON，不要输出 Markdown。字段必须稳定，缺失信息要明确标注。",
            },
            {"role": "user", "content": prompt},
        ],
    }
    async with httpx.AsyncClient(timeout=120) as client:
        response = await client.post(f"{settings.deepseek_base_url.rstrip('/')}/chat/completions", headers=headers, json=body)
        response.raise_for_status()
        text = response.json()["choices"][0]["message"]["content"]
    match = re.search(r"\{.*\}", text, re.S)
    if not match:
        return {}
    return json.loads(match.group(0))


async def call_deepseek_text(prompt: str, system_prompt: str) -> str:
    if settings.mock_mode:
        return ""
    headers = {"Authorization": f"Bearer {settings.deepseek_api_key}", "Content-Type": "application/json"}
    body = {
        "model": settings.deepseek_model,
        "temperature": 0.2,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": prompt},
        ],
    }
    async with httpx.AsyncClient(timeout=120) as client:
        response = await client.post(f"{settings.deepseek_base_url.rstrip('/')}/chat/completions", headers=headers, json=body)
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"]


def parse_deepseek_models(payload: dict[str, Any]) -> list[dict[str, str]]:
    models = []
    for item in payload.get("data", []):
        model_id = str(item.get("id") or "").strip()
        if not model_id:
            continue
        models.append({"id": model_id, "owned_by": str(item.get("owned_by") or "")})
    return models


async def list_deepseek_models() -> list[dict[str, str]]:
    if settings.mock_mode:
        return []
    headers = {"Authorization": f"Bearer {settings.deepseek_api_key}", "Content-Type": "application/json"}
    async with httpx.AsyncClient(timeout=30) as client:
        response = await client.get(f"{settings.deepseek_base_url.rstrip('/')}/models", headers=headers)
        response.raise_for_status()
        return parse_deepseek_models(response.json())


def save_analysis_result(db: Session, project: models.Project, payload: dict[str, Any], mode: str) -> models.ProjectReport:
    report = models.ProjectReport(
        project_id=project.id,
        report_type="project_analysis",
        content_json=json.dumps(payload, ensure_ascii=False),
        markdown=analysis_to_markdown(payload),
        model_name=settings.deepseek_model if mode == "deepseek" else "mock",
        mode=mode,
    )
    db.add(report)
    db.execute(delete(models.ProjectTask).where(models.ProjectTask.project_id == project.id))
    db.execute(delete(models.ProjectTimeline).where(models.ProjectTimeline.project_id == project.id))
    db.execute(delete(models.KnowledgeReference).where(models.KnowledgeReference.project_id == project.id))
    for item in payload.get("tasks", []):
        db.add(
            models.ProjectTask(
                project_id=project.id,
                task_name=item.get("task_name", ""),
                task_type=item.get("task_type", ""),
                priority=item.get("priority", "medium"),
                owner_role=item.get("owner_role", ""),
                estimated_days=int(item.get("estimated_days") or 1),
                dependencies=json.dumps(item.get("dependencies") or [], ensure_ascii=False),
                risk_level=item.get("risk_level", "medium"),
                status=item.get("status", "todo"),
                output_requirement=item.get("output_requirement", ""),
            )
        )
    for item in payload.get("timeline", []):
        db.add(
            models.ProjectTimeline(
                project_id=project.id,
                stage_name=item.get("stage_name", ""),
                start_day=int(item.get("start_day") or 1),
                end_day=int(item.get("end_day") or item.get("start_day") or 1),
                milestone=item.get("milestone", ""),
                dependencies=json.dumps(item.get("dependencies") or [], ensure_ascii=False),
                risk_note=item.get("risk_note", ""),
            )
        )
    for item in payload.get("knowledge_refs", []):
        db.add(
            models.KnowledgeReference(
                project_id=project.id,
                source_file=item.get("source_file", ""),
                source_path=item.get("source_path", ""),
                chunk_id=item.get("chunk_id", ""),
                quote=item.get("quote", ""),
                relevance_score=float(item.get("relevance_score") or 0),
            )
        )
    db.commit()
    db.refresh(report)
    return report


def ensure_project_sidecars(db: Session, project: models.Project) -> None:
    if project.tasks and project.timelines:
        return
    payload = mock_analysis_payload(project, project.files)
    save_analysis_result(db, project, payload, "mock")
    db.refresh(project)


def team_requirements_from_payload(payload: dict[str, Any]) -> dict[str, Any]:
    value = payload.get("team_requirements")
    if isinstance(value, dict):
        return value
    return {"total_headcount": 0, "roles": []}


def extract_md_metadata(text: str) -> tuple[str, list[str], list[str]]:
    heading = ""
    for line in text.splitlines():
        if line.startswith("#"):
            heading = line.lstrip("#").strip()
            break
    tags = sorted(set(re.findall(r"(?<!\w)#([\w\u4e00-\u9fff/-]+)", text)))
    links = sorted(set(re.findall(r"\[\[([^\]]+)\]\]", text)))
    return heading, tags, links


def chunk_text(text: str, path: str, tags: list[str], links: list[str]) -> list[dict[str, Any]]:
    if not text.strip():
        return []
    chunks = []
    current_heading = ""
    buffer = []
    for line in text.splitlines():
        if line.startswith("#") and buffer:
            content = "\n".join(buffer).strip()
            if content:
                chunks.append({"heading": current_heading, "content": content, "path": path, "tags": tags, "links": links})
            current_heading = line.lstrip("#").strip()
            buffer = [line]
        else:
            if line.startswith("#") and not current_heading:
                current_heading = line.lstrip("#").strip()
            buffer.append(line)
    content = "\n".join(buffer).strip()
    if content:
        chunks.append({"heading": current_heading, "content": content, "path": path, "tags": tags, "links": links})
    final = []
    for chunk in chunks:
        content = chunk["content"]
        if len(content) <= 1600:
            final.append(chunk)
            continue
        for i in range(0, len(content), 1400):
            final.append({**chunk, "content": content[i : i + 1600]})
    return final


def safe_browser_relative_path(filename: str) -> str:
    normalized = filename.replace("\\", "/").lstrip("/")
    parts = [part for part in PurePosixPath(normalized).parts if part not in {"", ".", ".."} and ":" not in part]
    return "/".join(parts) or "uploaded-file"


def index_knowledge_file(db: Session, source_path: Path, display_path: Optional[str] = None) -> dict[str, Any]:
    ext = source_path.suffix.lower()
    text, status = parse_document(source_path)
    heading, tags, links = extract_md_metadata(text) if ext == ".md" else ("", [], [])
    indexed_path = display_path or str(source_path)
    record = db.scalar(select(models.KnowledgeFile).where(models.KnowledgeFile.filepath == indexed_path))
    if not record:
        record = models.KnowledgeFile(filepath=indexed_path)
        db.add(record)
        db.flush()
    record.filename = Path(indexed_path).name
    record.filetype = ext.lstrip(".")
    record.filesize = source_path.stat().st_size
    record.title = heading or source_path.stem
    record.folder = str(Path(indexed_path).parent)
    db.execute(delete(models.KnowledgeChunk).where(models.KnowledgeChunk.file_id == record.id))
    db.execute(delete(models.KnowledgeTag).where(models.KnowledgeTag.file_id == record.id))
    db.execute(delete(models.KnowledgeLink).where(models.KnowledgeLink.file_id == record.id))
    for chunk in chunk_text(text, indexed_path, tags, links):
        db.add(
            models.KnowledgeChunk(
                file_id=record.id,
                heading=chunk["heading"],
                content=chunk["content"],
                path=chunk["path"],
                tags=json.dumps(chunk["tags"], ensure_ascii=False),
                links=json.dumps(chunk["links"], ensure_ascii=False),
            )
        )
    for tag, count in Counter(tags).items():
        db.add(models.KnowledgeTag(file_id=record.id, tag=tag, count=count))
    for link in links:
        db.add(models.KnowledgeLink(file_id=record.id, source_path=indexed_path, target=link))
    return {"filename": record.filename, "path": indexed_path, "filetype": record.filetype, "status": status}


def upsert_knowledge_markdown(db: Session, indexed_path: str, title: str, markdown: str, tags: list[str]) -> models.KnowledgeFile:
    record = db.scalar(select(models.KnowledgeFile).where(models.KnowledgeFile.filepath == indexed_path))
    if not record:
        record = models.KnowledgeFile(filepath=indexed_path)
        db.add(record)
        db.flush()
    record.filename = Path(indexed_path).name
    record.filetype = "md"
    record.filesize = len(markdown.encode("utf-8"))
    record.title = title
    record.folder = str(Path(indexed_path).parent)
    db.execute(delete(models.KnowledgeChunk).where(models.KnowledgeChunk.file_id == record.id))
    db.execute(delete(models.KnowledgeTag).where(models.KnowledgeTag.file_id == record.id))
    db.execute(delete(models.KnowledgeLink).where(models.KnowledgeLink.file_id == record.id))
    for chunk in chunk_text(markdown, indexed_path, tags, []):
        db.add(
            models.KnowledgeChunk(
                file_id=record.id,
                heading=chunk["heading"],
                content=chunk["content"],
                path=chunk["path"],
                tags=json.dumps(chunk["tags"], ensure_ascii=False),
                links=json.dumps(chunk["links"], ensure_ascii=False),
            )
        )
    for tag, count in Counter(tags).items():
        db.add(models.KnowledgeTag(file_id=record.id, tag=tag, count=count))
    return record


def scan_knowledge_directory(db: Session, directory: Path, clear_existing: bool = False) -> dict[str, Any]:
    if clear_existing:
        db.execute(delete(models.KnowledgeLink))
        db.execute(delete(models.KnowledgeTag))
        db.execute(delete(models.KnowledgeChunk))
        db.execute(delete(models.KnowledgeFile))
        db.commit()
    stats = Counter()
    recent = []
    for path in directory.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in SUPPORTED_KNOWLEDGE_EXTS:
            continue
        ext = path.suffix.lower()
        stats["total_files"] += 1
        stats[ext.lstrip(".") or "other"] += 1
        recent.append(index_knowledge_file(db, path))
    db.commit()
    return {"indexed_files": stats["total_files"], "stats": knowledge_stats(db), "skipped_files": [], "recent_files": recent[:20]}


def should_index_vault_path(path: Path, root: Path, include_sync_notes: bool = False) -> tuple[bool, str]:
    try:
        relative = path.relative_to(root)
    except ValueError:
        relative = path
    parts = set(relative.parts)
    if parts & DEFAULT_VAULT_EXCLUDE_DIRS:
        return False, "excluded_directory"
    if not include_sync_notes and "笔记同步助手" in parts:
        return False, "sync_notes_optional"
    if path.suffix.lower() not in SUPPORTED_KNOWLEDGE_EXTS:
        return False, "unsupported_type"
    if path.suffix.lower() in DEFAULT_VAULT_EXCLUDE_EXTS:
        return False, "excluded_type"
    try:
        if path.stat().st_size == 0:
            return False, "empty_file"
        if path.stat().st_size > DEFAULT_VAULT_MAX_FILE_SIZE:
            return False, "large_file"
    except OSError:
        return False, "stat_failed"
    return True, "ok"


def scan_vault_directory(
    db: Session,
    directory: Path,
    clear_existing: bool = False,
    include_sync_notes: bool = False,
) -> dict[str, Any]:
    if clear_existing:
        db.execute(delete(models.KnowledgeLink))
        db.execute(delete(models.KnowledgeTag))
        db.execute(delete(models.KnowledgeChunk))
        db.execute(delete(models.KnowledgeFile))
        db.commit()
    stats = Counter()
    skipped = Counter()
    recent = []
    priority_prefixes = ("wiki", "raw", "Obj-")
    for path in directory.rglob("*"):
        if not path.is_file():
            continue
        allowed, reason = should_index_vault_path(path, directory, include_sync_notes=include_sync_notes)
        if not allowed:
            skipped[reason] += 1
            continue
        try:
            rel = path.relative_to(directory)
            first = rel.parts[0] if rel.parts else ""
            priority = first == "wiki" or first == "raw" or first.startswith("Obj-")
        except ValueError:
            priority = False
        indexed = index_knowledge_file(db, path)
        indexed["priority_source"] = priority
        recent.append(indexed)
        stats["total_files"] += 1
        stats[path.suffix.lower().lstrip(".") or "other"] += 1
        if any(str(path).find(prefix) >= 0 for prefix in priority_prefixes):
            stats["priority_files"] += 1
    db.commit()
    return {
        "indexed_files": stats["total_files"],
        "stats": knowledge_stats(db),
        "skipped": dict(skipped),
        "skipped_files": [{"reason": reason, "count": count} for reason, count in skipped.items()],
        "recent_files": recent[:30],
        "filters": {
            "excluded_dirs": sorted(DEFAULT_VAULT_EXCLUDE_DIRS),
            "max_file_mb": DEFAULT_VAULT_MAX_FILE_SIZE // 1024 // 1024,
            "include_sync_notes": include_sync_notes,
        },
    }


def scan_vault_directory_with_progress(
    db: Session,
    directory: Path,
    clear_existing: bool = False,
    include_sync_notes: bool = False,
    progress: Optional[dict[str, Any]] = None,
) -> dict[str, Any]:
    if clear_existing:
        db.execute(delete(models.KnowledgeLink))
        db.execute(delete(models.KnowledgeTag))
        db.execute(delete(models.KnowledgeChunk))
        db.execute(delete(models.KnowledgeFile))
        db.commit()
    stats = Counter()
    skipped = Counter()
    recent = []
    paths = [path for path in directory.rglob("*") if path.is_file()]
    total = len(paths)
    if progress is not None:
        progress.update({"total_candidates": total, "processed": 0, "indexed_files": 0, "skipped_files": 0})
    priority_prefixes = ("wiki", "raw", "Obj-")
    for index, path in enumerate(paths, start=1):
        if progress is not None:
            progress.update({"processed": index, "current_file": str(path)})
        allowed, reason = should_index_vault_path(path, directory, include_sync_notes=include_sync_notes)
        if not allowed:
            skipped[reason] += 1
            if progress is not None:
                progress["skipped_files"] = sum(skipped.values())
            continue
        try:
            rel = path.relative_to(directory)
            first = rel.parts[0] if rel.parts else ""
            priority = first == "wiki" or first == "raw" or first.startswith("Obj-")
        except ValueError:
            priority = False
        indexed = index_knowledge_file(db, path)
        indexed["priority_source"] = priority
        recent.append(indexed)
        stats["total_files"] += 1
        stats[path.suffix.lower().lstrip(".") or "other"] += 1
        if any(str(path).find(prefix) >= 0 for prefix in priority_prefixes):
            stats["priority_files"] += 1
        if progress is not None:
            progress["indexed_files"] = stats["total_files"]
        if stats["total_files"] % 25 == 0:
            db.commit()
    db.commit()
    result = {
        "indexed_files": stats["total_files"],
        "stats": knowledge_stats(db),
        "skipped": dict(skipped),
        "skipped_files": [{"reason": reason, "count": count} for reason, count in skipped.items()],
        "recent_files": recent[:30],
        "filters": {
            "excluded_dirs": sorted(DEFAULT_VAULT_EXCLUDE_DIRS),
            "max_file_mb": DEFAULT_VAULT_MAX_FILE_SIZE // 1024 // 1024,
            "include_sync_notes": include_sync_notes,
        },
    }
    if progress is not None:
        progress.update({"processed": total, "current_file": "", "result": result})
    return result


def knowledge_overview(db: Session) -> dict[str, Any]:
    stats = knowledge_stats(db)
    files = list(db.scalars(select(models.KnowledgeFile).order_by(models.KnowledgeFile.updated_at.desc()).limit(80)))
    folder_counts = Counter(file.folder or "未分组" for file in files)
    type_counts = stats.get("filetype_distribution", {})
    recent = [
        {
            "filename": file.filename,
            "filepath": file.filepath,
            "filetype": file.filetype,
        }
        for file in files[:12]
    ]
    return {
        "stats": stats,
        "top_folders": [{"folder": folder, "count": count} for folder, count in folder_counts.most_common(10)],
        "type_counts": type_counts,
        "recent_files": recent,
    }


def is_knowledge_overview_question(question: str) -> bool:
    normalized = question.strip().replace("？", "").replace("?", "")
    return any(trigger in normalized for trigger in KNOWLEDGE_OVERVIEW_TRIGGERS)


def knowledge_overview_answer(db: Session) -> tuple[str, list[dict[str, Any]]]:
    overview = knowledge_overview(db)
    stats = overview["stats"]
    if not stats.get("total_files"):
        return "知识库目前还没有索引内容。请先索引本地路径或上传文件夹。", []
    type_counts = stats.get("filetype_distribution", {})
    type_text = "、".join(f"{key}: {value}" for key, value in sorted(type_counts.items())) or "暂无类型统计"
    folder_lines = "\n".join(f"- {item['folder']}：{item['count']} 个近期文件" for item in overview["top_folders"][:8])
    recent_lines = "\n".join(f"- {item['filename']}（{item['filetype']}）\n  {item['filepath']}" for item in overview["recent_files"][:8])
    answer = (
        "你的知识库当前已建立索引，可以从整体结构上这样理解：\n\n"
        f"- 总文件：{stats['total_files']} 个\n"
        f"- Markdown：{stats['markdown_files']} 个\n"
        f"- PDF / Word / Excel：{stats['pdf_docx_xlsx_files']} 个\n"
        f"- 图片：{stats['image_files']} 个\n"
        f"- 双链：{stats['link_count']} 条\n"
        f"- 文件类型分布：{type_text}\n\n"
        "近期资料主要分布在：\n"
        f"{folder_lines or '- 暂无分组'}\n\n"
        "最近索引的代表文件：\n"
        f"{recent_lines or '- 暂无文件'}\n\n"
        "你可以继续问更具体的问题，例如某个项目、某类规范、某个阶段任务或某份文件。"
    )
    refs = [
        {
            "chunk_id": "",
            "file_name": item["filename"],
            "file_path": item["filepath"],
            "quote": "知识库概览代表文件",
            "heading": "",
        }
        for item in overview["recent_files"][:8]
    ]
    return answer, refs


def knowledge_stats(db: Session) -> dict[str, Any]:
    files = list(db.scalars(select(models.KnowledgeFile)))
    tags = list(db.scalars(select(models.KnowledgeTag)))
    links_count = db.scalar(select(func.count(models.KnowledgeLink.id))) or 0
    type_counts = Counter(file.filetype for file in files)
    top_tags = Counter()
    for tag in tags:
        top_tags[tag.tag] += tag.count
    return {
        "total_files": len(files),
        "markdown_files": type_counts.get("md", 0),
        "pdf_docx_xlsx_files": type_counts.get("pdf", 0) + type_counts.get("docx", 0) + type_counts.get("xlsx", 0),
        "image_files": type_counts.get("png", 0) + type_counts.get("jpg", 0) + type_counts.get("jpeg", 0),
        "filetype_distribution": dict(type_counts),
        "top_tags": [{"tag": tag, "count": count} for tag, count in top_tags.most_common(20)],
        "link_count": links_count,
    }


def knowledge_tree(db: Session) -> list[dict[str, Any]]:
    files = list(db.scalars(select(models.KnowledgeFile).order_by(models.KnowledgeFile.filepath)))
    roots: dict[str, Any] = {}
    for file in files:
        parts = Path(file.filepath).parts
        current = roots
        for part in parts[:-1]:
            current = current.setdefault(part, {})
        current.setdefault("_files", []).append({"name": file.filename, "path": file.filepath, "filetype": file.filetype})

    def pack(name: str, value: Any) -> dict[str, Any]:
        children = [pack(k, v) for k, v in value.items() if k != "_files"]
        files_nodes = [{"name": f["name"], "type": "file", "path": f["path"], "filetype": f["filetype"]} for f in value.get("_files", [])]
        return {"name": name, "type": "folder", "children": children + files_nodes}

    return [pack(k, v) for k, v in roots.items()]


def list_knowledge_files(db: Session, q: str = "", limit: int = 100) -> dict[str, Any]:
    safe_limit = max(1, min(limit or 100, 500))
    query = select(models.KnowledgeFile).order_by(models.KnowledgeFile.updated_at.desc())
    if q.strip():
        keyword = f"%{q.strip()}%"
        query = query.where(or_(models.KnowledgeFile.filename.like(keyword), models.KnowledgeFile.filepath.like(keyword)))
    all_items = list(db.scalars(query))
    items = all_items[:safe_limit]
    return {
        "total": len(all_items),
        "items": [
            {
                "id": item.id,
                "filename": item.filename,
                "filepath": item.filepath,
                "filetype": item.filetype,
                "filesize": item.filesize,
                "updated_at": item.updated_at,
            }
            for item in items
        ],
    }


def search_knowledge(db: Session, question: str, limit: int = 6) -> list[models.KnowledgeChunk]:
    terms = [term for term in re.split(r"\s+", question.strip()) if term]
    if terms and settings.database_url.startswith("sqlite"):
        try:
            db.execute(text("CREATE VIRTUAL TABLE IF NOT EXISTS knowledge_chunks_fts USING fts5(chunk_id UNINDEXED, heading, content, path)"))
            db.execute(text("DELETE FROM knowledge_chunks_fts"))
            db.execute(
                text(
                    "INSERT INTO knowledge_chunks_fts(chunk_id, heading, content, path) "
                    "SELECT id, heading, content, path FROM knowledge_chunks"
                )
            )
            fts_query = " OR ".join(term.replace('"', "") for term in terms[:6])
            rows = list(
                db.execute(
                    text(
                        "SELECT chunk_id FROM knowledge_chunks_fts "
                        "WHERE knowledge_chunks_fts MATCH :q "
                        "ORDER BY bm25(knowledge_chunks_fts) LIMIT :limit"
                    ),
                    {"q": fts_query, "limit": limit},
                )
            )
            ids = [row[0] for row in rows]
            if ids:
                by_id = {chunk.id: chunk for chunk in db.scalars(select(models.KnowledgeChunk).where(models.KnowledgeChunk.id.in_(ids)))}
                ordered = [by_id[item_id] for item_id in ids if item_id in by_id]
                if ordered:
                    return ordered
        except Exception:
            db.rollback()
    query = select(models.KnowledgeChunk)
    if terms:
        clauses = []
        for term in terms[:6]:
            clauses.append(models.KnowledgeChunk.content.ilike(f"%{term}%"))
            clauses.append(models.KnowledgeChunk.heading.ilike(f"%{term}%"))
        query = query.where(or_(*clauses))
    chunks = list(db.scalars(query.limit(limit)))
    if chunks or not terms:
        return chunks
    return list(db.scalars(select(models.KnowledgeChunk).order_by(models.KnowledgeChunk.updated_at.desc()).limit(limit)))


def _markdown_list(items: list[str]) -> str:
    return "\n".join(f"- {item}" for item in items)


def _refs_from_chunks(chunks: list[models.KnowledgeChunk]) -> list[dict[str, Any]]:
    refs = []
    for chunk in chunks:
        refs.append(
            {
                "chunk_id": chunk.id,
                "source_file": Path(chunk.path).name,
                "source_path": chunk.path,
                "heading": chunk.heading,
                "quote": chunk.content[:420],
                "relevance_score": 0.82,
            }
        )
    return refs


def build_startup_analysis_payload(project: models.Project, chunks: list[models.KnowledgeChunk]) -> dict[str, Any]:
    refs = _refs_from_chunks(chunks)
    ref_names = [ref["source_file"] for ref in refs[:6]]
    description = project.description or "暂无项目描述，需由项目经理补充业主诉求、规模、阶段和交付目标。"
    technical_focus_cards = [
        {
            "title": "技术卡_日照计算",
            "dimension": "日照",
            "summary": "启动阶段应确认当地日照规则、分析口径、遮挡边界和是否已有日照复核成果。",
            "checkpoints": ["大寒日/冬至日口径", "被遮挡对象", "楼栋间距", "复核图纸版本"],
            "source_refs": ref_names,
            "manual_confirm": "需要人工确认当地规划部门采用的日照计算口径。",
        },
        {
            "title": "技术卡_退界要求",
            "dimension": "退界",
            "summary": "优先核对红线、道路、邻地、消防登高面及地下边界退距。",
            "checkpoints": ["用地红线", "道路退距", "邻地退距", "地下室边界", "消防登高面"],
            "source_refs": ref_names,
            "manual_confirm": "需要以规划条件或设计任务书为准。",
        },
        {
            "title": "技术卡_面积计算",
            "dimension": "面积",
            "summary": "明确计容/不计容、赠送面积、架空层、人防、地下室和配套用房计算方式。",
            "checkpoints": ["计容总建面", "可售面积", "地下室", "架空层", "配套用房", "人防面积"],
            "source_refs": ref_names,
            "manual_confirm": "需要和当地面积测绘/报规口径交叉确认。",
        },
        {
            "title": "技术卡_消防风险",
            "dimension": "消防",
            "summary": "强排前置检查消防车道、登高场地、间距、防火分区和地下车库组织。",
            "checkpoints": ["消防车道", "登高场地", "防火间距", "地下车库", "疏散组织"],
            "source_refs": ref_names,
            "manual_confirm": "需要专业负责人复核。",
        },
        {
            "title": "技术卡_规划条件",
            "dimension": "规划",
            "summary": "项目启动必须锁定容积率、限高、密度、绿地率、停车、配套和地块边界。",
            "checkpoints": ["容积率", "限高", "建筑密度", "绿地率", "停车配比", "配套要求"],
            "source_refs": ref_names,
            "manual_confirm": "缺规划条件时，不应进入最终任务拆解。",
        },
        {
            "title": "技术卡_报批风险",
            "dimension": "报批",
            "summary": "识别需要提前沟通的规划、消防、人防、面积、日照和专家会风险。",
            "checkpoints": ["规划沟通", "消防审查", "人防口径", "面积测算", "日照复核", "专家会"],
            "source_refs": ref_names,
            "manual_confirm": "需要项目经理形成待业主确认清单。",
        },
    ]
    task_breakdown = [
        {
            "task_name": "项目启动资料完整性检查",
            "task_type": "启动检查",
            "priority": "高",
            "owner_role": "项目经理",
            "estimated_days": 1,
            "dependencies": [],
            "risk_level": "中",
            "status": "todo",
            "output_requirement": "形成资料缺口表，确认规划条件、红线、任务书和历史参考资料。",
        },
        {
            "task_name": "历史项目技术复用提取",
            "task_type": "知识复用",
            "priority": "高",
            "owner_role": "AI资料管理员",
            "estimated_days": 1,
            "dependencies": ["项目启动资料完整性检查"],
            "risk_level": "中",
            "status": "todo",
            "output_requirement": "输出日照、退界、面积、消防、规划、报批六类技术卡。",
        },
        {
            "task_name": "启动会议程与业主确认清单",
            "task_type": "会议推进",
            "priority": "高",
            "owner_role": "项目经理",
            "estimated_days": 1,
            "dependencies": ["历史项目技术复用提取"],
            "risk_level": "低",
            "status": "todo",
            "output_requirement": "形成启动会议程、待确认问题和下一步责任人。",
        },
        {
            "task_name": "项目启动汇报PPT结构",
            "task_type": "汇报组织",
            "priority": "中",
            "owner_role": "AI汇报助理",
            "estimated_days": 2,
            "dependencies": ["启动会议程与业主确认清单"],
            "risk_level": "中",
            "status": "todo",
            "output_requirement": "输出面向业主的PPT目录和每页核心表达。",
        },
    ]
    meeting_agenda = [
        "确认项目背景、区位、规模、设计阶段和业主诉求。",
        "核对规划条件、红线、日照、退界、面积计算和消防/报批资料缺口。",
        "复用历史项目经验，明确哪些规则可参考、哪些必须重新确认。",
        "确认本轮交付物、责任人和时间节点。",
        "形成下次会议前的待办和业主决策事项。",
    ]
    ppt_outline = [
        {"page": 1, "title": "项目启动背景", "content": "项目基本信息、业主诉求、当前阶段。"},
        {"page": 2, "title": "资料完整度与缺口", "content": "已上传资料、缺失资料、需确认来源。"},
        {"page": 3, "title": "历史项目技术复用", "content": "日照、退界、面积、消防、规划、报批六类重点。"},
        {"page": 4, "title": "启动阶段任务拆解", "content": "任务、负责人、优先级、交付物。"},
        {"page": 5, "title": "风险与下一步", "content": "未决问题、业主决策事项、下一次会议。"},
    ]
    risk_list = [
        "规划条件、红线或任务书缺失会导致任务拆解偏差。",
        "日照、退界、面积计算口径必须以当地正式要求为准。",
        "历史项目只能作为复用参考，不能替代本项目审批依据。",
        "若会议结论不回写任务看板，项目经理后续追踪会断链。",
    ]
    open_questions = [
        "业主本轮最关心的是进度、产品定位、成本还是报批风险？",
        "是否已有正式规划条件、红线和设计任务书？",
        "本项目是否需要先做强排可行性或日照快速复核？",
        "启动会后哪些事项需要业主书面确认？",
    ]
    mindmap_json = {
        "title": "项目启动分析",
        "nodes": [
            {
                "id": "project",
                "label": project.name,
                "children": [
                    {
                        "id": "technical",
                        "label": "技术重点",
                        "children": [{"id": card["dimension"], "label": card["dimension"]} for card in technical_focus_cards],
                    },
                    {
                        "id": "tasks",
                        "label": "任务拆解",
                        "children": [{"id": item["task_name"], "label": item["task_name"]} for item in task_breakdown],
                    },
                    {
                        "id": "meeting",
                        "label": "启动会",
                        "children": [{"id": str(idx), "label": item} for idx, item in enumerate(meeting_agenda, start=1)],
                    },
                    {
                        "id": "ppt",
                        "label": "PPT结构",
                        "children": [{"id": str(item["page"]), "label": item["title"]} for item in ppt_outline],
                    },
                ],
            }
        ],
    }
    return {
        "mode": "mock" if settings.mock_mode else "local_workflow",
        "project_summary": {
            "name": project.name,
            "city": project.city,
            "project_type": project.project_type,
            "phase": project.phase,
            "description": description,
            "knowledge_refs_count": len(refs),
            "summary": f"{project.name} 当前进入项目启动分析，重点是把资料缺口、历史技术复用和会议推进转成可执行任务。",
        },
        "technical_focus_cards": technical_focus_cards,
        "task_breakdown": task_breakdown,
        "meeting_agenda": meeting_agenda,
        "ppt_outline": ppt_outline,
        "risk_list": risk_list,
        "open_questions": open_questions,
        "mindmap_json": mindmap_json,
        "source_refs": refs,
    }


def pending_analysis_files(project: models.Project) -> list[models.ProjectFile]:
    return [file for file in project.files if (file.analysis_status or "pending") == "pending"]


def _project_file_context(project: models.Project, limit: int = 8, files: Optional[list[models.ProjectFile]] = None) -> list[dict[str, str]]:
    contexts = []
    source_files = files if files is not None else project.files
    for file in source_files[:limit]:
        text = (file.parsed_text or "").strip()
        status = file.parse_status or "pending"
        if not text:
            text = "文件尚未解析出文本，请在分析中标记为资料未解析或仅可作为文件存在性参考。"
        contexts.append(
            {
                "file_id": file.id,
                "filename": file.filename,
                "filetype": file.filetype,
                "parse_status": status,
                "analysis_status": file.analysis_status or "pending",
                "text": text[:5000],
            }
        )
    return contexts


def build_startup_analysis_prompt(project: models.Project, chunks: list[models.KnowledgeChunk], files: Optional[list[models.ProjectFile]] = None) -> str:
    file_context = _project_file_context(project, files=files)
    source_refs = _refs_from_chunks(chunks)
    expected_schema = {
        "project_summary": {
            "name": "项目名称",
            "city": "城市",
            "project_type": "项目类型",
            "phase": "阶段",
            "description": "项目描述",
            "knowledge_refs_count": 0,
            "summary": "基于上传文件和知识库资料形成的真实启动分析摘要",
        },
        "technical_focus_cards": [
            {
                "title": "技术卡标题",
                "dimension": "日照/退界/面积/消防/规划/报批等维度",
                "summary": "结合上传文件的判断",
                "checkpoints": ["需要复核的具体事项"],
                "source_refs": ["引用的文件名或知识库来源"],
                "manual_confirm": "需要人工确认的内容或风险等级",
            }
        ],
        "task_breakdown": [
            {
                "task_name": "任务名称",
                "task_type": "任务类型",
                "priority": "高/中/低",
                "owner_role": "负责人角色",
                "estimated_days": 1,
                "dependencies": [],
                "risk_level": "高/中/低",
                "status": "todo",
                "output_requirement": "交付物要求",
            }
        ],
        "meeting_agenda": ["启动会要讨论的问题"],
        "ppt_outline": [{"page": 1, "title": "页面标题", "content": "页面核心内容"}],
        "risk_list": ["风险点"],
        "open_questions": ["待确认问题"],
        "mindmap_json": {"title": "项目启动分析", "nodes": []},
        "source_refs": [
            {
                "chunk_id": "来源ID",
                "source_file": "来源文件",
                "source_path": "来源路径",
                "heading": "标题",
                "quote": "引用片段",
                "relevance_score": 0.8,
            }
        ],
    }
    return (
        "请基于项目基本信息、用户上传并解析的文件内容、以及知识库检索片段，生成建筑设计项目启动分析 JSON。\n"
        "必须优先分析上传文件里的真实内容；如果资料不足，请明确指出缺口，不要编造不存在的条件。\n"
        "只输出 JSON，不要输出 Markdown，不要解释。\n\n"
        "项目基本信息：\n"
        f"{json.dumps({'name': project.name, 'city': project.city, 'project_type': project.project_type, 'phase': project.phase, 'description': project.description}, ensure_ascii=False, indent=2)}\n\n"
        "本次待分析文件内容：\n"
        f"{json.dumps(file_context, ensure_ascii=False, indent=2)}\n\n"
        "知识库检索片段：\n"
        f"{json.dumps(source_refs, ensure_ascii=False, indent=2)}\n\n"
        "输出 JSON Schema 示例：\n"
        f"{json.dumps(expected_schema, ensure_ascii=False, indent=2)}"
    )


def normalize_startup_analysis_payload(payload: dict[str, Any], fallback: dict[str, Any]) -> dict[str, Any]:
    merged = dict(payload or {})
    for key in [
        "project_summary",
        "technical_focus_cards",
        "task_breakdown",
        "meeting_agenda",
        "ppt_outline",
        "risk_list",
        "open_questions",
        "mindmap_json",
        "source_refs",
    ]:
        if not merged.get(key):
            merged[key] = fallback.get(key)
    if not isinstance(merged.get("project_summary"), dict):
        merged["project_summary"] = fallback.get("project_summary", {})
    if not isinstance(merged.get("mindmap_json"), dict):
        merged["mindmap_json"] = fallback.get("mindmap_json", {"title": "项目启动分析", "nodes": []})
    for key in ["technical_focus_cards", "task_breakdown", "meeting_agenda", "ppt_outline", "risk_list", "open_questions", "source_refs"]:
        if not isinstance(merged.get(key), list):
            merged[key] = fallback.get(key, [])
    merged["mode"] = merged.get("mode") or "deepseek"
    return merged


def startup_analysis_to_markdown(payload: dict[str, Any]) -> str:
    summary = payload.get("project_summary", {})
    lines = [
        "# 项目启动分析",
        "",
        f"## 项目摘要",
        f"- 项目：{summary.get('name', '')}",
        f"- 阶段：{summary.get('phase', '')}",
        f"- 结论：{summary.get('summary', '')}",
        "",
        "## 技术重点",
    ]
    for card in payload.get("technical_focus_cards", []):
        lines.append(f"- **{card.get('dimension')}**：{card.get('summary')}")
    lines.extend(["", "## 任务拆解"])
    for task in payload.get("task_breakdown", []):
        lines.append(f"- {task.get('task_name')}（{task.get('owner_role')} / {task.get('priority')}）")
    lines.extend(["", "## 启动会议程"])
    lines.extend(f"{idx}. {item}" for idx, item in enumerate(payload.get("meeting_agenda", []), start=1))
    lines.extend(["", "## PPT结构"])
    for item in payload.get("ppt_outline", []):
        lines.append(f"{item.get('page')}. {item.get('title')}：{item.get('content')}")
    lines.extend(["", "## 风险与未决问题"])
    lines.extend(f"- {item}" for item in payload.get("risk_list", []))
    lines.extend(f"- 待确认：{item}" for item in payload.get("open_questions", []))
    return "\n".join(lines)


def _upsert_startup_skill_card(
    db: Session,
    project: models.Project,
    card_type: str,
    title: str,
    markdown: str,
    output: dict[str, Any],
) -> models.SkillCard:
    card = db.scalar(
        select(models.SkillCard).where(models.SkillCard.project_id == project.id, models.SkillCard.card_type == card_type)
    )
    if not card:
        card = models.SkillCard(project_id=project.id, card_type=card_type)
        db.add(card)
    card.title = title
    card.status = "succeeded"
    card.input_data = json.dumps({"source": "startup_analysis"}, ensure_ascii=False)
    card.output_data = json.dumps(output, ensure_ascii=False)
    card.input_json = json.dumps({"source": "startup_analysis"}, ensure_ascii=False)
    card.output_json = json.dumps(output, ensure_ascii=False)
    card.markdown = markdown
    card.source = "startup_analysis"
    card.created_by = "system"
    card.completed_at = datetime.utcnow()
    return card


def _candidate_markdown(title: str, body: str, sources: list[str]) -> str:
    lines = [
        "---",
        "type: obsidian_candidate",
        f'title: "{title}"',
        "status: pending_review",
        f'updated: "{datetime.utcnow().date().isoformat()}"',
        "---",
        "",
        f"# {title}",
        "",
        body,
        "",
        "## 候选来源",
    ]
    lines.extend(f"- {source}" for source in sources[:8])
    return "\n".join(lines)


def startup_analysis_deposit_markdown(project: models.Project, payload: dict[str, Any], report_id: str) -> str:
    summary = payload.get("project_summary", {})
    lines = [
        "---",
        "type: project_deposit",
        f'project_id: "{project.id}"',
        f'project_name: "{project.name}"',
        f'city: "{project.city}"',
        f'phase: "{project.phase}"',
        f'report_id: "{report_id}"',
        f'updated: "{datetime.utcnow().date().isoformat()}"',
        "---",
        "",
        f"# {project.name} 项目分析沉淀",
        "",
        "## 项目摘要",
        summary.get("summary") or project.description or "暂无摘要。",
        "",
        "## 可复用技术重点",
    ]
    for card in payload.get("technical_focus_cards", []):
        lines.append(f"### {card.get('title') or card.get('dimension') or '技术重点'}")
        lines.append(card.get("summary") or "")
        checkpoints = card.get("checkpoints") or []
        if checkpoints:
            lines.append("")
            lines.append("复核要点：")
            lines.extend(f"- {item}" for item in checkpoints)
        source_refs = card.get("source_refs") or []
        if source_refs:
            lines.append("")
            lines.append("参考来源：")
            lines.extend(f"- {item}" for item in source_refs[:6])
        lines.append("")
    lines.append("## 任务拆解经验")
    for task in payload.get("task_breakdown", []):
        lines.append(
            f"- {task.get('task_name', '')}：{task.get('owner_role', '')}，优先级 {task.get('priority', '')}，交付物：{task.get('output_requirement', '')}"
        )
    lines.extend(["", "## 风险与待确认问题"])
    lines.extend(f"- 风险：{item}" for item in payload.get("risk_list", []))
    lines.extend(f"- 待确认：{item}" for item in payload.get("open_questions", []))
    lines.extend(["", "## 原始引用来源"])
    for ref in payload.get("source_refs", [])[:12]:
        source = ref.get("source_path") or ref.get("source_file") or "未命名来源"
        quote = ref.get("quote") or ""
        lines.append(f"- {source}：{quote[:160]}")
    return "\n".join(lines).strip() + "\n"


def deposit_startup_analysis_to_knowledge(
    db: Session,
    project: models.Project,
    payload: dict[str, Any],
    report_id: str,
) -> models.KnowledgeFile:
    title = f"{project.name} 项目分析沉淀"
    indexed_path = f"project-deposits/{project.id}/startup-analysis.md"
    tags = ["项目沉淀", "启动分析", project.name, project.city, project.project_type, project.phase]
    markdown = startup_analysis_deposit_markdown(project, payload, report_id)
    record = upsert_knowledge_markdown(db, indexed_path, title, markdown, [tag for tag in tags if tag])
    db.flush()
    return record


def meeting_summary_without_transcript_excerpt(summary: str) -> str:
    marker = "\n## 原始转写摘录\n"
    if marker not in summary:
        return summary
    return summary.split(marker, 1)[0].rstrip()


def meeting_summary_deposit_markdown(project: models.Project, meeting: models.Meeting) -> str:
    lines = [
        "---",
        "type: meeting_summary",
        f"project_id: {project.id}",
        f"project: {project.name}",
        f"meeting_id: {meeting.id}",
        f"meeting_title: {meeting.title}",
        f"status: {meeting.status}",
        "---",
        "",
        f"# {project.name} - {meeting.title} AI会议纪要",
        "",
        "## AI纪要",
        meeting_summary_without_transcript_excerpt(meeting.summary or "暂无AI纪要。"),
    ]
    actions = meeting.next_actions_json or ""
    if actions:
        lines.extend(["", "## 待办", actions])
    if (meeting.transcript or "").strip():
        lines.extend(
            [
                "",
                "## 原始转写",
                f"完整转写已独立保存：project-deposits/{project.id}/meeting-transcripts/{meeting.id}.md",
            ]
        )
    return "\n".join(lines).strip() + "\n"


def deposit_meeting_summary_to_knowledge(
    db: Session,
    project: models.Project,
    meeting: models.Meeting,
) -> models.KnowledgeFile:
    title = f"{project.name} {meeting.title} AI会议纪要"
    indexed_path = f"project-deposits/{project.id}/meetings/{meeting.id}.md"
    tags = ["项目沉淀", "会议纪要", "AI纪要", project.name, meeting.title, project.city, project.project_type, project.phase]
    markdown = meeting_summary_deposit_markdown(project, meeting)
    record = upsert_knowledge_markdown(db, indexed_path, title, markdown, [tag for tag in tags if tag])
    db.flush()
    return record


def deposit_meeting_transcript_to_knowledge(
    db: Session,
    project: models.Project,
    meeting: models.Meeting,
) -> Optional[models.KnowledgeFile]:
    transcript = (meeting.transcript or "").strip()
    if not transcript:
        return None
    title = f"{project.name} {meeting.title} 原始转写"
    indexed_path = f"project-deposits/{project.id}/meeting-transcripts/{meeting.id}.md"
    markdown = "\n".join(
        [
            "---",
            "type: meeting_transcript",
            f"project_id: {project.id}",
            f"meeting_id: {meeting.id}",
            f"meeting_title: {meeting.title}",
            "---",
            "",
            f"# {title}",
            "",
            transcript,
        ]
    ).strip() + "\n"
    tags = ["项目沉淀", "会议原始转写", project.name, meeting.title]
    record = upsert_knowledge_markdown(db, indexed_path, title, markdown, tags)
    db.flush()
    return record


def save_startup_analysis(
    db: Session,
    project: models.Project,
    payload: dict[str, Any],
    analyzed_files: Optional[list[models.ProjectFile]] = None,
) -> models.ProjectReport:
    db.execute(
        delete(models.ProjectReport).where(
            models.ProjectReport.project_id == project.id,
            models.ProjectReport.report_type == "startup_analysis",
        )
    )
    db.execute(delete(models.KnowledgeReference).where(models.KnowledgeReference.project_id == project.id))
    report = models.ProjectReport(
        project_id=project.id,
        report_type="startup_analysis",
        content_json=json.dumps(payload, ensure_ascii=False),
        markdown=startup_analysis_to_markdown(payload),
        model_name=settings.deepseek_model if not settings.mock_mode else "mock",
        mode=payload.get("mode", "mock"),
    )
    db.add(report)
    for ref in payload.get("source_refs", [])[:12]:
        db.add(
            models.KnowledgeReference(
                project_id=project.id,
                source_file=ref.get("source_file", ""),
                source_path=ref.get("source_path", ""),
                chunk_id=ref.get("chunk_id", ""),
                quote=ref.get("quote", ""),
                relevance_score=float(ref.get("relevance_score") or 0),
            )
        )
    existing_task_names = {task.task_name for task in project.tasks}
    for item in payload.get("task_breakdown", []):
        if item["task_name"] in existing_task_names:
            continue
        db.add(
            models.ProjectTask(
                project_id=project.id,
                task_name=item["task_name"],
                task_type=item["task_type"],
                priority=item["priority"],
                owner_role=item["owner_role"],
                estimated_days=int(item["estimated_days"]),
                dependencies=json.dumps(item.get("dependencies", []), ensure_ascii=False),
                risk_level=item["risk_level"],
                status=item["status"],
                output_requirement=item["output_requirement"],
            )
        )
    agenda_markdown = "# 项目启动会议程\n\n" + "\n".join(f"{idx}. {item}" for idx, item in enumerate(payload["meeting_agenda"], start=1))
    if not project.meetings:
        db.add(
            models.Meeting(
                project_id=project.id,
                title=f"{project.name} 项目启动会",
                agenda=agenda_markdown,
                status="scheduled",
                mindmap_json=json.dumps(payload["mindmap_json"], ensure_ascii=False),
                next_actions_json=json.dumps(payload["task_breakdown"], ensure_ascii=False),
            )
        )
    _upsert_startup_skill_card(
        db,
        project,
        "technical_focus",
        "技术重点卡",
        "## 技术重点\n" + "\n".join(f"- **{card['dimension']}**：{card['summary']}" for card in payload["technical_focus_cards"]),
        {"cards": payload["technical_focus_cards"]},
    )
    _upsert_startup_skill_card(
        db,
        project,
        "task_breakdown",
        "任务拆解卡",
        "## 启动任务\n" + "\n".join(f"- {task['task_name']}（{task['owner_role']}）" for task in payload["task_breakdown"]),
        {"items": payload["task_breakdown"]},
    )
    _upsert_startup_skill_card(
        db,
        project,
        "meeting_agenda",
        "会议议程卡",
        agenda_markdown,
        {"items": payload["meeting_agenda"]},
    )
    _upsert_startup_skill_card(
        db,
        project,
        "ppt_outline",
        "PPT结构卡",
        "## PPT结构\n" + "\n".join(f"{item['page']}. {item['title']}：{item['content']}" for item in payload["ppt_outline"]),
        {"slides": payload["ppt_outline"]},
    )
    sources = [ref.get("source_path", "") for ref in payload.get("source_refs", [])]
    db.execute(
        delete(models.KnowledgeItem).where(
            models.KnowledgeItem.project_id == project.id,
            models.KnowledgeItem.item_type == "obsidian_candidate",
        )
    )
    for card in payload.get("technical_focus_cards", []):
        db.add(
            models.KnowledgeItem(
                project_id=project.id,
                source_file=card["title"],
                item_type="obsidian_candidate",
                summary=card["summary"],
                tags=json.dumps(["技术卡", card["dimension"], "候选"], ensure_ascii=False),
                content=_candidate_markdown(card["title"], card["summary"], sources),
            )
        )
    batch_id = uuid4().hex
    for file in analyzed_files or []:
        file.analysis_status = "analyzed"
        file.analysis_batch_id = batch_id
        file.analyzed_at = datetime.utcnow()
    db.flush()
    deposit_startup_analysis_to_knowledge(db, project, payload, report.id)
    db.commit()
    db.refresh(report)
    return report


async def run_startup_analysis(db: Session, project: models.Project) -> dict[str, Any]:
    question = " ".join(
        [
            project.name,
            project.city,
            project.project_type,
            project.phase,
            project.description,
            "日照 退界 面积计算 消防 规划条件 报批 强排 会议纪要",
        ]
    )
    chunks = search_knowledge(db, question, limit=10)
    analysis_files = pending_analysis_files(project)
    fallback_payload = build_startup_analysis_payload(project, chunks)
    if analysis_files:
        fallback_payload["analysis_scope"] = {
            "mode": "pending_files_only",
            "file_count": len(analysis_files),
            "filenames": [file.filename for file in analysis_files],
        }
    payload = fallback_payload
    if not settings.mock_mode:
        try:
            deepseek_payload = await call_deepseek_json(build_startup_analysis_prompt(project, chunks, files=analysis_files))
            if deepseek_payload:
                payload = normalize_startup_analysis_payload(deepseek_payload, fallback_payload)
                payload["mode"] = "deepseek"
            else:
                raise RuntimeError("DeepSeek 未返回有效结构化分析")
        except Exception as exc:
            raise RuntimeError(f"DeepSeek 真实分析失败：{exc}") from exc
    report = save_startup_analysis(db, project, payload, analyzed_files=analysis_files)
    return {"report": report, **payload}


def classify_project_execution_instruction(project: models.Project, instruction: str) -> str:
    text = instruction.strip().lower()
    compact = re.sub(r"[\s，。！？!?,.、~～]+", "", text)
    if compact in {"你好", "您好", "hello", "hi", "嗨", "哈喽", "在吗", "在不在"}:
        return "greeting"

    project_terms = {
        project.name,
        project.city,
        project.project_type,
        project.phase,
        "项目",
        "资料",
        "任务",
        "会议",
        "纪要",
        "执行",
        "设计",
        "方案",
        "强排",
        "总图",
        "日照",
        "退界",
        "消防",
        "面积",
        "容积率",
        "报批",
        "规划",
        "楼栋",
        "户型",
        "配套",
        "风险",
        "知识库",
    }
    if any(term and str(term).lower() in text for term in project_terms):
        return "project"
    return "unrelated"


def build_project_execution_greeting(project: models.Project) -> str:
    return f"你好，我在。你可以直接问我和“{project.name}”有关的问题，我会按当前项目资料、任务、会议纪要和知识库来回答。"


def build_project_execution_refusal(project: models.Project) -> str:
    return f"这个问题和当前项目“{project.name}”没有直接关系，我先不展开回答。请发和项目资料、任务、会议、设计风险或推进安排有关的问题。"


def build_project_execution_prompt(project: models.Project, instruction: str, chunks: list[models.KnowledgeChunk]) -> str:
    latest_report = next((report for report in sorted(project.reports, key=lambda item: item.created_at, reverse=True)), None)
    context = {
        "project": {
            "name": project.name,
            "city": project.city,
            "project_type": project.project_type,
            "phase": project.phase,
            "description": project.description,
            "status": project.status,
        },
        "latest_report": (latest_report.markdown[:3000] if latest_report and latest_report.markdown else ""),
        "tasks": [
            {
                "task_name": task.task_name,
                "owner_role": task.owner_role,
                "priority": task.priority,
                "risk_level": task.risk_level,
                "status": task.status,
                "output_requirement": task.output_requirement,
            }
            for task in project.tasks[:20]
        ],
        "meetings": [
            {
                "title": meeting.title,
                "agenda": (meeting.agenda or "")[:800],
                "summary": (meeting.summary or "")[:1200],
                "status": meeting.status,
            }
            for meeting in project.meetings[:8]
        ],
        "knowledge_references": [
            {
                "source_file": ref.source_file,
                "source_path": ref.source_path,
                "quote": ref.quote[:500],
            }
            for ref in project.knowledge_references[:10]
        ],
        "retrieved_knowledge": [
            {
                "path": chunk.path,
                "heading": chunk.heading,
                "content": chunk.content[:900],
            }
            for chunk in chunks
        ],
    }
    return (
        "你是建筑设计项目执行助手。请基于项目上下文、本地知识库检索结果和用户指令执行工作。\n"
        "要求：\n"
        "1. 直接回应用户发来的具体问题，不要套固定项目模板，不要答非所问。\n"
        "2. 必须优先使用项目上下文和知识库内容，不要编造来源。\n"
        "3. 如果用户问题与当前项目无关，只能礼貌说明无法在项目执行台回答无关内容。\n"
        "4. 普通打招呼可以自然回应，不要强行输出项目分析。\n"
        "5. 输出要服务项目推进；只有在用户问题需要时，才包含判断、拆解、风险和下一步建议。\n"
        "6. 如果适合转成任务，请列出任务名、负责人角色、优先级和交付物。\n"
        "7. 如使用了知识库内容，可以在答案末尾简短列出来源路径；没有使用来源时不要写引用或来源说明。\n\n"
        f"用户指令：{instruction}\n\n"
        "项目执行上下文：\n"
        f"{json.dumps(context, ensure_ascii=False, indent=2)}"
    )


async def run_project_execution(db: Session, project: models.Project, instruction: str) -> models.AgentRun:
    instruction_type = classify_project_execution_instruction(project, instruction)
    chunks: list[models.KnowledgeChunk] = []
    prompt = ""
    if instruction_type == "greeting":
        mode = "greeting"
        answer = build_project_execution_greeting(project)
    elif instruction_type == "unrelated":
        mode = "refused"
        answer = build_project_execution_refusal(project)
    else:
        query = " ".join([project.name, project.city, project.project_type, project.phase, instruction])
        chunks = search_knowledge(db, query, limit=8)
        prompt = build_project_execution_prompt(project, instruction, chunks)
        mode = "mock" if settings.mock_mode else "deepseek"
        try:
            if settings.mock_mode:
                answer = (
                    f"【Mock模式】已读取 {project.name} 的项目上下文，并检索到 {len(chunks)} 条知识库片段。\n\n"
                    f"针对“{instruction}”，建议先围绕资料缺口、技术风险、任务责任人和交付物四类事项拆解。\n\n"
                    "## 下一步建议\n"
                    "- 明确本轮要解决的关键技术问题。\n"
                    "- 对照知识库历史项目经验形成复核清单。\n"
                    "- 将结论转为任务或沉淀为项目知识。"
                )
            else:
                answer = await call_deepseek_text(
                    prompt=prompt,
                    system_prompt="你是建筑设计项目执行助手。只基于用户当前问题、项目上下文和本地知识库回答；无关内容礼貌拒绝；普通问候自然回应；输出中文 Markdown。",
                )
        except Exception as exc:
            mode = "deepseek_error"
            answer = (
                "DeepSeek 调用失败，已回退为本地执行摘要。\n\n"
                f"- 指令：{instruction}\n"
                f"- 已检索知识片段：{len(chunks)} 条\n"
                f"- 错误类型：{exc.__class__.__name__}\n\n"
                "建议稍后重试，或先基于当前项目任务和会议纪要手动推进。"
            )
    output = {
        "mode": mode,
        "answer": answer,
        "references": [
            {
                "chunk_id": chunk.id,
                "source_path": chunk.path,
                "heading": chunk.heading,
                "quote": chunk.content[:420],
            }
            for chunk in chunks
        ],
    }
    run = models.AgentRun(
        project_id=project.id,
        agent_id="project-execution",
        input_context=json.dumps({"instruction": instruction, "instruction_type": instruction_type, "prompt": prompt}, ensure_ascii=False),
        output_json=json.dumps(output, ensure_ascii=False),
        status="succeeded" if mode != "deepseek_error" else "failed",
    )
    db.add(run)
    db.commit()
    db.refresh(run)
    return run


def build_team_plan(db: Session, project: models.Project) -> models.TeamPlan:
    employees = list(db.scalars(select(models.DigitalEmployee)))
    tasks = project.tasks
    role_tasks: dict[str, list[str]] = defaultdict(list)
    for task in tasks:
        role_tasks[task.owner_role or "AI项目经理"].append(task.task_name)
    roles = []
    for employee in employees:
        matched = role_tasks.get(employee.name) or role_tasks.get(employee.role) or []
        if matched or employee.name in {"AI项目经理", "AI资料管理员", "AI汇报助理"}:
            roles.append(
                {
                    "name": employee.name,
                    "role": employee.role,
                    "recommended_count": 1,
                    "tasks": matched[:6],
                    "skills": json.loads(employee.skills or "[]"),
                    "intensity": "高" if matched else "中",
                    "risk_note": "需要人工确认任务边界和交付标准。",
                }
            )
    if not roles:
        roles = [
            {
                "name": "AI项目经理",
                "role": "项目负责人 / PM",
                "recommended_count": 1,
                "tasks": [],
                "skills": ["任务拆解"],
                "intensity": "中",
                "risk_note": "Mock 规则生成。",
            }
        ]
    plan = models.TeamPlan(
        project_id=project.id,
        recommended_roles=json.dumps(roles, ensure_ascii=False),
        staffing_summary=f"建议 {len(roles)} 类数字员工参与；需要根据真实任务继续校准投入强度。",
    )
    db.add(plan)
    db.commit()
    db.refresh(plan)
    return plan


def default_meeting_agenda(project: models.Project) -> str:
    return "\n".join(
        [
            f"# {project.name} 项目启动会议程",
            "1. 确认项目背景、业主诉求与当前设计阶段。",
            "2. 核对规划条件、红线、日照、退界、面积计算等技术资料缺口。",
            "3. 复盘历史相似项目可复用经验和风险清单。",
            "4. 明确本轮交付物、节点时间和责任人。",
            "5. 确认下一次业主沟通需要决策的问题。",
        ]
    )


def _format_meeting_summary_markdown(project_name: str, payload: dict[str, Any]) -> str:
    summary = str(payload.get("summary") or "").strip()
    core_items = [str(item).strip() for item in payload.get("core_items", []) if str(item).strip()]
    demand_translation = payload.get("demand_translation", []) or []
    decisions = [str(item).strip() for item in payload.get("decisions", []) if str(item).strip()]
    risks = [str(item).strip() for item in payload.get("risks", []) if str(item).strip()]
    next_steps = [str(item).strip() for item in payload.get("next_steps", []) if str(item).strip()]
    broadcast_script = str(payload.get("broadcast_script") or "").strip()
    lines = [f"# {project_name} 五段式会议纪要", "", "## 1. 纪要内容", summary or "未能从真实转写中提取明确摘要。"]
    if core_items:
        lines.extend(["", "## 2. 核心事项", *[f"- {item}" for item in core_items]])
    if demand_translation:
        lines.extend(["", "## 3. 甲方诉求转译（内部研判）"])
        for item in demand_translation:
            if isinstance(item, dict):
                raw = str(item.get("raw") or item.get("source") or "原话待确认").strip()
                time = str(item.get("time") or "时间点待确认").strip()
                translation = str(item.get("translation") or item.get("meaning") or "转译待复核").strip()
                action = str(item.get("design_response") or item.get("action") or "").strip()
                lines.append(f"- 原话：{raw}（{time}）")
                lines.append(f"  转译：{translation}")
                if action:
                    lines.append(f"  设计回应：{action}")
            else:
                lines.append(f"- {item}")
    if decisions:
        lines.extend(["", "## 4. 会议决议", *[f"- {item}" for item in decisions]])
    if risks:
        lines.extend(["", "## 风险与问题", *[f"- {item}" for item in risks]])
    if next_steps:
        lines.extend(["", "## 5. 待办事项", *[f"- {item}" for item in next_steps]])
    if broadcast_script:
        lines.extend(["", "## 语音播报稿", broadcast_script])
    return "\n".join(lines).strip()


def _fallback_real_transcript_summary(project_name: str, transcript: str) -> tuple[str, list[dict[str, str]]]:
    excerpt = transcript.strip()[:1800]
    summary = (
        f"# {project_name} 会议纪要\n\n"
        "## 摘要\n"
        "以下内容基于腾讯会议真实转写生成。当前模型未返回结构化纪要，先保留原始转写摘要供复核。\n\n"
        "## 原始转写摘录\n"
        f"{excerpt}"
    )
    return summary, []


async def summarize_meeting(db: Session, meeting: models.Meeting) -> models.Meeting:
    project = db.get(models.Project, meeting.project_id)
    project_name = project.name if project else "当前项目"
    transcript = (meeting.transcript or "").strip()
    if not transcript:
        raise ValueError("请先同步腾讯会议真实转写，再生成AI纪要")

    prompt = (
        "请只基于以下腾讯会议真实转写生成五段式会议纪要，不要补充转写中没有的信息。\n"
        "输出 JSON，字段：summary 字符串；core_items 字符串数组；"
        "demand_translation 数组，每项包含 raw、time、translation、design_response，"
        "仅转译真实出现的甲方模糊诉求，例如“不够高级/没气势/不像某风格”；"
        "decisions 字符串数组；risks 字符串数组；next_steps 字符串数组；"
        "todos 数组，每项包含 title、owner、status；broadcast_script 字符串。"
        "如果无法判断负责人，owner 写“待确认”；status 默认 todo。\n\n"
        "种子词典：不够高级→材料质感/比例/留白/入口仪式感；没气势→体量轮廓/轴线/界面展开；"
        "不像宋式→屋面/檐口/色彩/尺度语汇；不够松弛→空间节奏/界面密度/景观留白。\n\n"
        f"项目：{project_name}\n"
        f"会议标题：{meeting.title}\n"
        f"真实转写：\n{transcript[:12000]}"
    )
    actions: list[dict[str, str]] = []
    try:
        payload = await call_deepseek_json(prompt)
    except Exception:
        payload = {}
    if payload:
        meeting.summary = _format_meeting_summary_markdown(project_name, payload)
        for item in payload.get("todos", []):
            if isinstance(item, dict):
                title = str(item.get("title") or "").strip()
                if title:
                    actions.append(
                        {
                            "title": title,
                            "owner": str(item.get("owner") or "待确认").strip() or "待确认",
                            "status": str(item.get("status") or "todo").strip() or "todo",
                        }
                    )
    else:
        meeting.summary, actions = _fallback_real_transcript_summary(project_name, transcript)

    meeting.mindmap_json = "{}"
    meeting.next_actions_json = json.dumps(actions, ensure_ascii=False)
    meeting.status = "summarized"
    if project:
        deposit_meeting_transcript_to_knowledge(db, project, meeting)
        deposit_meeting_summary_to_knowledge(db, project, meeting)
    db.commit()

    if project:
        for item in actions:
            existing = db.scalar(
                select(models.ProjectTask).where(
                    models.ProjectTask.project_id == project.id,
                    models.ProjectTask.source_type == "meeting",
                    models.ProjectTask.source_id == meeting.id,
                    models.ProjectTask.task_name == item["title"],
                )
            )
            if existing:
                continue
            db.add(
                models.ProjectTask(
                    project_id=project.id,
                    task_name=item["title"],
                    task_type="会议待办",
                    priority="高",
                    owner_role=item["owner"],
                    estimated_days=2,
                    dependencies="[]",
                    risk_level="中",
                    status="todo",
                    output_requirement="来自会议纪要自动生成，需要项目经理确认。",
                    source_type="meeting",
                    source_id=meeting.id,
                )
            )
        db.commit()
    db.refresh(meeting)
    return meeting


def _project_context_pack(project: models.Project) -> dict[str, Any]:
    files = list(project.files[:12])
    meetings = list(project.meetings[:6])
    tasks = list(project.tasks[:10])
    reports = list(project.reports[:4])
    refs = list(project.knowledge_references[:8])
    return {
        "project": {
            "id": project.id,
            "name": project.name,
            "city": project.city,
            "type": project.project_type,
            "phase": project.phase,
            "status": project.status,
            "description": project.description,
        },
        "files": [
            {
                "name": file.filename,
                "type": file.filetype,
                "parse_status": file.parse_status,
                "summary": (file.parsed_text or "")[:600],
            }
            for file in files
        ],
        "meetings": [
            {
                "title": meeting.title,
                "summary": (meeting.summary or "")[:800],
                "next_actions": meeting.next_actions_json,
            }
            for meeting in meetings
        ],
        "tasks": [
            {
                "name": task.task_name,
                "status": task.status,
                "priority": task.priority,
                "risk": task.risk_level,
                "owner": task.owner_role,
            }
            for task in tasks
        ],
        "reports": [
            {
                "type": report.report_type,
                "markdown": (report.markdown or "")[:1200],
            }
            for report in reports
        ],
        "knowledge_references": [
            {
                "source": ref.source_file,
                "quote": ref.quote,
                "score": ref.relevance_score,
            }
            for ref in refs
        ],
        "counts": {
            "files": len(project.files),
            "meetings": len(project.meetings),
            "tasks": len(project.tasks),
            "knowledge_refs": len(project.knowledge_references),
            "skill_cards": len(project.skill_cards),
        },
    }


def _skill_source_refs(chunks: list[models.KnowledgeChunk], project: models.Project) -> list[dict[str, str]]:
    refs = _refs_from_chunks(chunks) if chunks else []
    if refs:
        return [
            {
                "source_file": str(ref.get("source_file") or ref.get("source_path") or "知识库"),
                "quote": str(ref.get("quote") or "")[:240],
            }
            for ref in refs
        ]
    return [
        {"source_file": ref.source_file, "quote": (ref.quote or "")[:240]}
        for ref in project.knowledge_references[:6]
    ]


def _fallback_skill_output(skill: SkillDefinition, project: models.Project, prompt: str, sources: list[dict[str, str]]) -> dict[str, Any]:
    project_name = project.name or "当前项目"
    common_sources = sources or [{"source_file": "当前项目上下文", "quote": "尚未检索到明确知识库来源。"}]
    if skill.id == "brief_interpretation":
        return {
            "explicit_goals": ["确认项目任务书中的建设目标、阶段成果和汇报对象"],
            "implicit_goals": ["补齐甲方真实诉求、审查设计矛盾与边界条件"],
            "design_conflicts": ["资料完整性与方案判断之间仍有缺口"],
            "entry_points": ["先完成资料归档与任务书解读，再形成汇报主线"],
            "sources": common_sources,
        }
    if skill.id == "task_breakdown":
        return {
            "tasks": [
                {"title": "补齐项目基础资料", "owner": "项目经理", "priority": "高", "deliverable": "资料缺口清单"},
                {"title": "复核技术边界", "owner": "技术负责人", "priority": "高", "deliverable": "日照/退界/指标风险表"},
                {"title": "搭建汇报结构", "owner": "方案主创", "priority": "中", "deliverable": "PPT大纲"},
            ],
            "next": "可将任务写回任务看板后分派。",
        }
    if skill.id == "technical_focus":
        return {
            "cards": [
                {"dimension": "日照", "risk": "中", "checkpoints": ["日照标准", "遮挡关系", "测算口径"]},
                {"dimension": "退界", "risk": "中", "checkpoints": ["红线", "绿线", "消防登高面"]},
                {"dimension": "面积", "risk": "中", "checkpoints": ["容积率", "计容口径", "地库边界"]},
                {"dimension": "消防", "risk": "高", "checkpoints": ["消防车道", "登高场地", "防火分区"]},
            ],
            "sources": common_sources,
        }
    if skill.id == "meeting_minutes":
        return {
            "summary": f"{project_name} 本次会议围绕项目推进、资料补齐和下一步成果进行讨论。",
            "core_items": ["确认下一阶段汇报目标", "补齐技术边界资料", "形成可执行待办"],
            "demand_translation": [
                {
                    "raw": prompt[:160] or "待补充会议原话",
                    "time": "待确认",
                    "translation": "需要转化为材料、比例、界面、体量或汇报逻辑上的明确设计动作。",
                    "internal_only": True,
                }
            ],
            "decisions": ["先完成资料归档与技术边界复核，再进入汇报结构深化。"],
            "todos": [{"title": "整理会议待办并分配责任人", "owner": "项目经理", "status": "todo"}],
            "broadcast_script": "本次会议确认了资料补齐、技术边界复核和下一步汇报准备三项重点，请项目团队优先处理。",
            "sources": common_sources,
        }
    if skill.id == "ppt_outline":
        return {
            "title": f"{project_name} 方案汇报框架",
            "slides": [
                {"page": 1, "title": "封面", "content": "项目名称、阶段、汇报对象"},
                {"page": 2, "title": "项目背景", "content": "城市、区位、用地与甲方目标"},
                {"page": 3, "title": "核心问题", "content": "设计矛盾、技术边界、风险点"},
                {"page": 4, "title": "设计主线", "content": "概念、空间策略、产品策略"},
                {"page": 5, "title": "方案展开", "content": "总图、流线、户型/功能、立面方向"},
                {"page": 6, "title": "下一步", "content": "待补资料、关键节点与责任分工"},
            ],
            "key_messages": ["结论前置", "风险透明", "策略可执行"],
            "missing_assets": ["关键图纸", "技术指标", "参考案例"],
            "sources": common_sources,
        }
    if skill.id == "concept_copy":
        return {
            "concept_title": f"{project_name} 的场景更新与秩序重构",
            "narrative": "从场地真实约束和甲方诉求出发，建立清晰的空间秩序、识别性界面与可落地的产品表达。",
            "strategies": ["强化入口界面", "控制体量比例", "建立连续公共空间", "让材料与尺度回应项目定位"],
            "presentation_copy": "本方案以清晰的空间秩序回应场地约束，以克制而有识别度的建筑表达建立项目记忆点。",
            "sources": common_sources,
        }
    if skill.id == "competitor_analysis":
        return {
            "comparables": ["历史项目/案例库待检索后自动回填"],
            "transferable_strategies": ["入口仪式感", "材料层级", "展示界面", "汇报叙事结构"],
            "limits": ["不同城市、甲方、容积率与报批环境下不可直接套用。"],
            "risks": ["若缺少真实历史项目索引，竞品分析会退化为方法建议。"],
            "sources": common_sources,
        }
    if skill.id == "reference_image_classification":
        return {
            "image_type": "待分类参考图",
            "style_tags": ["现代", "克制", "展示面"],
            "material_tags": ["石材", "金属", "玻璃"],
            "reuse_points": ["入口界面", "立面比例", "材料层级"],
            "next_prompt": "可继续生成基于当前项目的生图提示词。",
        }
    if skill.id == "image_prompt":
        return {
            "positive_prompt": f"{project_name}，建筑方案前期意向图，结合城市语境、项目阶段与甲方诉求，现代克制，高品质材料，清晰体量比例，真实建筑摄影感",
            "negative_prompt": "低清晰度，过度奇幻，结构不合理，文字水印，畸形透视",
            "style_tags": ["建筑摄影", "方案意向", "高品质", "克制"],
            "camera": "eye-level architectural photography, 35mm lens",
            "usage": "用于方案前期风格探索与汇报意向沟通",
            "sources": common_sources,
        }
    if skill.id == "scheme_review":
        return {
            "strengths": ["项目已有资料可支撑基础判断"],
            "risks": ["技术边界、甲方诉求和成果缺口仍需复核"],
            "missing_info": ["完整任务书", "最新会议结论", "指标表", "关键图纸"],
            "next_actions": ["先补齐资料，再做方案评审定稿"],
            "sources": common_sources,
        }
    return {"result": "已生成基础成果。", "sources": common_sources}


def _markdown_from_output(skill: SkillDefinition, output: dict[str, Any], prompt: str) -> str:
    lines = [f"# {skill.name}", ""]
    if prompt:
        lines.extend(["## 用户需求", prompt, ""])
    for key, value in output.items():
        title = {
            "summary": "纪要内容",
            "core_items": "核心事项",
            "demand_translation": "甲方诉求转译",
            "decisions": "会议决议",
            "todos": "待办事项",
            "broadcast_script": "播报稿",
            "slides": "PPT框架",
            "tasks": "任务清单",
            "sources": "来源",
        }.get(key, key)
        lines.append(f"## {title}")
        if isinstance(value, list):
            for item in value:
                if isinstance(item, dict):
                    lines.append("- " + "；".join(f"{k}: {v}" for k, v in item.items() if k != "sources"))
                else:
                    lines.append(f"- {item}")
        elif isinstance(value, dict):
            for item_key, item_value in value.items():
                lines.append(f"- {item_key}: {item_value}")
        else:
            lines.append(str(value))
        lines.append("")
    return "\n".join(lines).strip()


def _normalize_skill_id(card_type: str) -> str:
    aliases = {
        "ppt_structure": "ppt_outline",
        "tech_points": "technical_focus",
        "image_generation": "ai_image_generation",
        "ai_image": "ai_image_generation",
    }
    return aliases.get(card_type, card_type)


async def _deepseek_skill_output(skill: SkillDefinition, project: models.Project, prompt: str, context: dict[str, Any], sources: list[dict[str, str]]) -> dict[str, Any]:
    if settings.mock_mode:
        return _fallback_skill_output(skill, project, prompt, sources)
    schema_hint = "、".join(skill.output_schema)
    instruction = (
        "你是建筑方案前期的项目 AI 执行器。必须围绕当前项目回答，避免通用空话。\n"
        "判断、研判、转译、方案、风险和复用类内容要优先基于项目上下文与知识库来源；不要编造来源。\n"
        "请输出 JSON，不要输出 Markdown。字段尽量贴合以下 schema："
        f"{schema_hint}。\n"
    )
    payload = {
        "skill": {"id": skill.id, "name": skill.name, "description": skill.description},
        "user_request": prompt,
        "project_context": context,
        "retrieved_sources": sources,
    }
    try:
        raw = await call_deepseek_text(
            prompt=json.dumps(payload, ensure_ascii=False, default=str),
            system_prompt=instruction,
        )
        match = re.search(r"\{.*\}", raw, re.S)
        if match:
            parsed = json.loads(match.group(0))
            if isinstance(parsed, dict):
                if sources and "sources" not in parsed:
                    parsed["sources"] = sources
                return parsed
    except Exception:
        pass
    return _fallback_skill_output(skill, project, prompt, sources)


async def _generate_image_asset(project: models.Project, prompt: str, output: dict[str, Any]) -> dict[str, Any]:
    image_prompt = (
        output.get("positive_prompt")
        or output.get("prompt")
        or prompt
        or f"{project.name} 建筑方案前期意向图"
    )
    result: dict[str, Any] = {
        "provider": settings.image_provider,
        "model": settings.image_model,
        "prompt": image_prompt,
        "image_paths": [],
        "status": "not_configured",
        "message": "图片生成服务未配置。",
    }
    if not settings.image_configured:
        return result

    target_dir = project_upload_dir(project.id) / "generated-images"
    target_dir.mkdir(parents=True, exist_ok=True)
    headers = {"Authorization": f"Bearer {settings.image_api_key}"}
    body = {
        "model": settings.image_model,
        "prompt": image_prompt,
        "n": 1,
        "size": "1024x1024",
    }
    try:
        async with httpx.AsyncClient(timeout=120) as client:
            response = await client.post(f"{settings.image_base_url.rstrip('/')}/images/generations", headers=headers, json=body)
            response.raise_for_status()
            data = response.json()
            first = (data.get("data") or [{}])[0]
            image_bytes = b""
            if first.get("b64_json"):
                image_bytes = base64.b64decode(first["b64_json"])
            elif first.get("url"):
                image_response = await client.get(first["url"])
                image_response.raise_for_status()
                image_bytes = image_response.content
            if not image_bytes:
                raise RuntimeError("图片服务未返回可保存的图片数据")
            file_path = target_dir / f"ai-image-{datetime.now().strftime('%Y%m%d-%H%M%S')}-{uuid4().hex[:6]}.png"
            file_path.write_bytes(image_bytes)
            result.update({"image_paths": [str(file_path)], "status": "succeeded", "message": "图片已生成并保存到项目成果目录。"})
    except Exception as exc:
        result.update({"status": "failed", "message": f"图片生成失败：{exc}"})
    return result


async def _execute_skill(db: Session, project: models.Project, skill: SkillDefinition, prompt: str, intent_meta: Optional[dict[str, Any]] = None) -> models.SkillCard:
    context = _project_context_pack(project)
    query = " ".join([prompt, project.name, project.city, project.project_type, project.phase]).strip()
    chunks = search_knowledge(db, query, limit=8) if (skill.retrieval_required and query) else []
    sources = _skill_source_refs(chunks, project)
    output = await _deepseek_skill_output(skill, project, prompt, context, sources)

    if skill.id == "ai_image_generation":
        prompt_skill = SKILL_BY_ID.get("image_prompt")
        prompt_output = await _deepseek_skill_output(prompt_skill or skill, project, prompt, context, sources)
        image_result = await _generate_image_asset(project, prompt, prompt_output)
        output = {**prompt_output, **image_result, "source_context": sources}

    markdown = _markdown_from_output(skill, output, prompt)
    input_payload = {
        "prompt": prompt,
        "project_id": project.id,
        "skill": skill.id,
        "intent": intent_meta or {},
        "context_counts": context.get("counts", {}),
    }
    card = models.SkillCard(
        project_id=project.id,
        card_type=skill.id,
        title=skill.name,
        status="succeeded" if output.get("status") not in {"failed", "not_configured"} else output.get("status", "succeeded"),
        input_json=json.dumps(input_payload, ensure_ascii=False, default=str),
        output_json=json.dumps(output, ensure_ascii=False, default=str),
        markdown=markdown,
        source="huashu" if skill.id == "ai_image_generation" else ("mock" if settings.mock_mode else "deepseek"),
        completed_at=datetime.utcnow(),
    )
    db.add(card)
    db.commit()
    db.refresh(card)
    return card


def run_skill_card(db: Session, project: models.Project, card_type: str, prompt: str = "") -> models.SkillCard:
    skill_id = _normalize_skill_id(card_type)
    skill = SKILL_BY_ID.get(skill_id, SKILL_BY_ID["task_breakdown"])
    sources = _skill_source_refs([], project)
    output = _fallback_skill_output(skill, project, prompt, sources)
    markdown = _markdown_from_output(skill, output, prompt)
    card = models.SkillCard(
        project_id=project.id,
        card_type=skill.id,
        title=skill.name,
        status="succeeded",
        input_json=json.dumps({"prompt": prompt, "project_id": project.id, "skill": skill.id}, ensure_ascii=False),
        output_json=json.dumps(output, ensure_ascii=False, default=str),
        markdown=markdown,
        source="local_skill",
        completed_at=datetime.utcnow(),
    )
    db.add(card)
    db.commit()
    db.refresh(card)
    return card


async def run_agent_chat(db: Session, project: models.Project, message: str, requested_skill: str = "") -> dict[str, Any]:
    if requested_skill:
        skill = SKILL_BY_ID.get(_normalize_skill_id(requested_skill), SKILL_BY_ID["task_breakdown"])
        confidence = 1.0
        reason = "用户直接指定技能。"
    else:
        skill, confidence, reason = match_skill_by_keywords(message)
        if confidence < 0.5 and not settings.mock_mode:
            try:
                choices = list_builtin_skills()
                raw = await call_deepseek_text(
                    prompt=json.dumps({"message": message, "skills": choices}, ensure_ascii=False),
                    system_prompt=(
                        "请为建筑设计 AI 代理选择最合适的 skill。只输出 JSON："
                        '{"intent":"skill_id","confidence":0.0到1.0,"reason":"原因"}。'
                    ),
                )
                match = re.search(r"\{.*\}", raw, re.S)
                if match:
                    payload = json.loads(match.group(0))
                    candidate = SKILL_BY_ID.get(_normalize_skill_id(str(payload.get("intent") or "")))
                    if candidate:
                        skill = candidate
                        confidence = float(payload.get("confidence") or confidence)
                        reason = str(payload.get("reason") or reason)
            except Exception:
                pass
    card = await _execute_skill(db, project, skill, message, {"confidence": confidence, "reason": reason})
    return {
        "intent": skill.id,
        "confidence": confidence,
        "reason": reason,
        "selected_skill": {
            "id": skill.id,
            "name": skill.name,
            "description": skill.description,
            "downstream": list(skill.downstream),
        },
        "card": card,
        "context": _project_context_pack(project).get("counts", {}),
        "available_skills": list_builtin_skills(),
    }
