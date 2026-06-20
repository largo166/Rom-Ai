"""core.py — 基础工具、常量、员工初始化、文件操作"""
import hashlib
import shutil
from pathlib import Path
from typing import Any, Optional
from uuid import uuid4

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from sqlalchemy import delete, select
from sqlalchemy.orm import Session

from app import models
from app.config import settings
from app.json_safety import safe_json_dump

# ──────────────────────────── 常量 ────────────────────────────

SUPPORTED_PROJECT_EXTS = {".pdf", ".docx", ".xlsx", ".txt", ".md", ".pptx", ".png", ".jpg", ".jpeg"}
SUPPORTED_INBOX_EXTS = SUPPORTED_PROJECT_EXTS | {".zip", ".rar", ".7z", ".dwg", ".skp", ".rvt", ".csv", ".webp", ".mp3", ".m4a", ".wav", ".mp4"}
SUPPORTED_KNOWLEDGE_EXTS = {".md", ".txt", ".pdf", ".docx", ".xlsx", ".pptx", ".png", ".jpg", ".jpeg"}
MAX_BROWSER_UPLOAD_SIZE = 25 * 1024 * 1024
DEFAULT_VAULT_EXCLUDE_DIRS = {
    ".git",
    ".obsidian",
    ".venv",
    ".venv-funasr",
    "__pycache__",
    "node_modules",
    "dist",
    ".cache",
}
DEFAULT_VAULT_EXCLUDE_EXTS = {
    ".mp4",
    ".mov",
    ".avi",
    ".zip",
    ".rar",
    ".7z",
    ".psd",
    ".dwg",
    ".bak",
}
DEFAULT_VAULT_MAX_FILE_SIZE = 30 * 1024 * 1024
KNOWLEDGE_OVERVIEW_TRIGGERS = ("知识库里有什么", "知识库有什么", "有什么资料", "资料概览", "知识库概览", "总结知识库")

DIGITAL_EMPLOYEES = [
    ("AI项目经理", "项目负责人 / PM", ["任务拆解", "进度协调", "风险跟踪"], "PM"),
    ("AI总图助理", "总图负责人", ["强排", "流线", "指标复核"], "MP"),
    ("AI户型助理", "户型负责人", ["产品定位", "户型推演", "得房率"], "UN"),
    ("AI立面助理", "立面负责人", ["风格研究", "材料策略", "比例控制"], "FC"),
    ("AI景观助理", "景观接口", ["景观节点", "归家动线", "界面协同"], "LS"),
    ("AI室内助理", "室内接口", ["大堂", "会所", "精装氛围"], "IN"),
    ("AI汇报助理", "汇报助理", ["PPT目录", "汇报文案", "图像提示词"], "RP"),
    ("AI资料管理员", "资料管理员", ["文件整理", "知识库引用", "版本记录"], "KB"),
]

TEAM_MEMBERS = [
    ("项目经理", "项目负责人", ["业主沟通", "会议推进", "任务协调"], 45),
    ("主创负责人", "主创设计师", ["概念判断", "方案把控", "评审"], 55),
    ("总图负责人", "总图负责人", ["强排", "指标复核", "退界"], 50),
    ("立面负责人", "立面负责人", ["风格研究", "材料策略", "比例控制"], 35),
    ("汇报负责人", "汇报助理", ["PPT结构", "文本组织", "成果整合"], 40),
]


# ──────────────────────────── 员工初始化 ────────────────────────────

def ensure_digital_employees(db: Session) -> None:
    existing_names = set(db.scalars(select(models.DigitalEmployee.name)))
    target_names = {item[0] for item in DIGITAL_EMPLOYEES}
    if target_names.issubset(existing_names):
        return
    db.execute(delete(models.DigitalEmployee))
    for name, role, skills, avatar in DIGITAL_EMPLOYEES:
        db.add(
            models.DigitalEmployee(
                name=name,
                role=role,
                skills=safe_json_dump(skills, field_name="skills"),
                avatar=avatar,
                status="available",
                workload=20,
            )
        )
    db.commit()


def ensure_team_members(db: Session) -> None:
    existing_names = set(db.scalars(select(models.TeamMember.name)))
    if existing_names:
        return
    for name, role, skills, workload in TEAM_MEMBERS:
        db.add(
            models.TeamMember(
                name=name,
                role=role,
                skills=safe_json_dump(skills, field_name="skills"),
                status="available",
                workload=workload,
            )
        )
    db.commit()


# ──────────────────────────── 文档解析 ────────────────────────────

def parse_document(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    try:
        if ext in {".png", ".jpg", ".jpeg"}:
            return "", "saved_no_ocr"
        if ext in {".txt", ".md"}:
            return path.read_text(encoding="utf-8", errors="ignore"), "parsed"
        if ext == ".pdf":
            reader = PdfReader(str(path))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
            return text, "parsed"
        if ext == ".docx":
            doc = Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
            return text, "parsed"
        if ext == ".xlsx":
            wb = load_workbook(str(path), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    values = [str(value) for value in row if value is not None]
                    if values:
                        rows.append(" | ".join(values))
            return "\n".join(rows), "parsed"
        if ext == ".pptx":
            try:
                from pptx import Presentation
            except Exception as exc:
                return f"解析失败：缺少 python-pptx（{exc}）", "failed"
            prs = Presentation(str(path))
            slides = []
            for idx, slide in enumerate(prs.slides, start=1):
                slides.append(f"# Slide {idx}")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        slides.append(text)
            return "\n".join(slides), "parsed"
        return "", "unsupported"
    except Exception as exc:
        return f"解析失败：{exc}", "failed"


# ──────────────────────────── 路径工具 ────────────────────────────

def project_upload_dir(project_id: str) -> Path:
    path = settings.upload_root_path / "projects" / project_id
    path.mkdir(parents=True, exist_ok=True)
    return path


def project_managed_dir(project_id: str, upload_root: Optional[Path] = None) -> Path:
    root = upload_root or settings.upload_root_path
    return root / "projects" / project_id


def inbox_upload_dir() -> Path:
    path = settings.upload_root_path / "inbox"
    path.mkdir(parents=True, exist_ok=True)
    return path


def delete_project_library(db: Session, project_id: str, upload_root: Optional[Path] = None) -> dict[str, Any]:
    project = db.get(models.Project, project_id)
    if not project:
        return {"deleted": False, "deleted_project_id": project_id, "deleted_files": 0}
    managed_dir = project_managed_dir(project_id, upload_root)
    deleted_files = len(project.files)
    inbox_items = list(db.scalars(select(models.InboxItem).where(models.InboxItem.project_id == project_id)))
    for item in inbox_items:
        item.project_id = ""
        item.archive_path = ""
        item.status = "待确认"
        item.archive_group = "待确认"
        item.needs_review = True
    db.delete(project)
    db.commit()
    if managed_dir.exists() and managed_dir.is_dir():
        shutil.rmtree(managed_dir, ignore_errors=True)
    return {"deleted": True, "deleted_project_id": project_id, "deleted_files": deleted_files}


def _unique_path(directory: Path, filename: str) -> Path:
    directory.mkdir(parents=True, exist_ok=True)
    safe_name = Path(filename).name or "uploaded-file"
    candidate = directory / safe_name
    if not candidate.exists():
        return candidate
    stem = candidate.stem
    suffix = candidate.suffix
    for index in range(2, 1000):
        next_candidate = directory / f"{stem}_{index}{suffix}"
        if not next_candidate.exists():
            return next_candidate
    return directory / f"{stem}_{uuid4().hex[:8]}{suffix}"


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()
